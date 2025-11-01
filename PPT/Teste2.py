import pandas as pd
from pptx import Presentation
import os
import re
import sys
from datetime import datetime

def criar_apresentacao_sinistro(numero_sinistro=None, excel_file=None, template_path=None):
    """
    Cria uma apresentação PowerPoint baseada em um template e dados de uma planilha Excel.
    
    Args:
        numero_sinistro: Número do sinistro para filtrar (opcional)
        excel_file: Caminho para o arquivo Excel (opcional)
        template_path: Caminho para o template PPT (opcional)
    
    Returns:
        str: Caminho do arquivo PPT gerado ou None em caso de erro
    """
    # Definir caminhos padrão se não forem fornecidos
    if excel_file is None:
        excel_file = 'SINISTROS_FILTRADO.xlsx'
    
    if template_path is None:
        template_path = 'Template.pptx'
    
    # Carregar os dados do Excel
    try:
        df = pd.read_excel(excel_file)
    except Exception as e:
        print(f"Erro ao ler o arquivo Excel: {e}")
        return None
    
    # Verificar se há dados
    if df.empty:
        print("A planilha Excel está vazia.")
        return None
    
    # Se um número de sinistro for fornecido, filtrar os dados
    if numero_sinistro is not None:
        if "Nº Reguladora" in df.columns:
            df = df[df["Nº Reguladora"] == numero_sinistro]
            if df.empty:
                print(f"Sinistro {numero_sinistro} não encontrado na planilha.")
                return None
    
    # Extrair os dados do sinistro (primeira linha da planilha)
    sinistro_data = df.iloc[0]
    
    # Obter o número do sinistro dos dados
    if numero_sinistro is None:
        if "Nº Reguladora" in sinistro_data:
            numero_sinistro = sinistro_data["Nº Reguladora"]
        else:
            numero_sinistro = "Desconhecido"
    
    # Definir o caminho de saída
    output_path = f'C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\Sinistro_{numero_sinistro}_Novo.pptx'
    
    # Formatar o valor do prejuízo
    prejuizo = sinistro_data.get("Prejuizo Apurado", "N/A")
    if pd.notna(prejuizo) and isinstance(prejuizo, (int, float)):
        prejuizo_formatado = f"R$ {prejuizo:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        prejuizo_formatado = str(prejuizo)
    
    # Formatar a data do sinistro
    data_sinistro = sinistro_data.get("Data do Sinistro", "N/A")
    if isinstance(data_sinistro, datetime):
        data_formatada = data_sinistro.strftime("%d/%m/%Y")
    else:
        data_formatada = str(data_sinistro)
    
    # Carregar o template PPT
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Erro ao carregar o template PPT: {e}")
        return None
    
    # Mapear os dados do sinistro para os campos do template
    sinistro_info = {
        'Número do Sinistro': str(numero_sinistro),
        'Status': str(sinistro_data.get("Observação", "")) if pd.notna(sinistro_data.get("Observação", "")) else "",
        'Data': data_formatada,
        'Causa': str(sinistro_data.get("Causa Final", "")) if pd.notna(sinistro_data.get("Causa Final", "")) else "",
        'Transportador': str(sinistro_data.get("Transportador", "")) if pd.notna(sinistro_data.get("Transportador", "")) else "",
        'UF Origem': str(sinistro_data.get("UF - Origem", "")) if pd.notna(sinistro_data.get("UF - Origem", "")) else "",
        'Cidade Origem': str(sinistro_data.get("Cidade Origem", "")) if pd.notna(sinistro_data.get("Cidade Origem", "")) else "",
        'UF Destino': str(sinistro_data.get("UF - Destino", "")) if pd.notna(sinistro_data.get("UF - Destino", "")) else "",
        'Cidade Destino': str(sinistro_data.get("Cidade - Destino", "")) if pd.notna(sinistro_data.get("Cidade - Destino", "")) else "",
        'Prejuízo Apurado': prejuizo_formatado,
        'Descrição': str(sinistro_data.get("Ação", "")) if pd.notna(sinistro_data.get("Ação", "")) else ""
    }
    
    # Função para substituir texto mantendo a formatação
    def replace_text_in_shape(shape, replacements):
        if not hasattr(shape, "text_frame"):
            return
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, value)
    
    # Atualizar os slides com os dados do sinistro
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                replace_text_in_shape(shape, sinistro_info)
            
            # Processar tabelas, se houver
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_shape(cell, sinistro_info)
    
    # Salvar a apresentação
    try:
        prs.save(output_path)
        print(f"Apresentação gerada com sucesso: {output_path}")
        return output_path
    except Exception as e:
        print(f"Erro ao salvar a apresentação: {e}")
        return None

if __name__ == "__main__":
    # Se um número de sinistro for fornecido como argumento, use-o
    if len(sys.argv) > 1:
        try:
            numero_sinistro = int(sys.argv[1])
            criar_apresentacao_sinistro(numero_sinistro)
        except ValueError:
            print("Erro: O número do sinistro deve ser um valor numérico.")
    else:
        # Caso contrário, use os dados da primeira linha da planilha
        criar_apresentacao_sinistro()