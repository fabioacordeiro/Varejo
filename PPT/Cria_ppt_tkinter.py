# COMO USAR
# 1) Python 3.10+
# 2) Instale dependências:
#    pip install python-pptx pandas openpyxl numpy
# 3) Execute:
#    python Cria_ppt_Fabio.py
#
# OBSERVAÇÕES IMPORTANTES
# - Este GUI usa Tkinter/ttk padrão (sem libs externas de tema).
# - Ele mantém a lógica do script (ler Excel, filtrar sinistro, escolher template A/R/V,
#   substituir textos no PPT e salvar).
# - Ele remove input() e print(): tudo vai para um log na tela e também para o console.
# - Geração em thread para não travar a interface.
# Created by: Fabio A Cordeiro
# Date: 24/02/2026 11:35

from __future__ import annotations

import os
import re
import sys
import threading
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from queue import Queue, Empty
from typing import Optional, Dict, Any, List

import numpy as np
import pandas as pd
from pptx import Presentation
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter import scrolledtext
import logging


# ============================================================
# CONFIG / DEFAULTS
# ============================================================

DEFAULT_EXCEL = Path(r"C:\\Fabio\\Desenvolvimento\\Varejo\Sinistro\\SINISTROS1.xlsx")
DEFAULT_OUTPUT_DIR = Path(r"C:\\Fabio\\Desenvolvimento\\Varejo\\PPT")

# Ajuste aqui caso seus templates fiquem em outra pasta
TEMPLATE_PATHS = {
    "A": Path(r"C:\\Fabio\Desenvolvimento\\Varejo\\PPT\\TemplateA.pptx"),
    "R": Path(r"C:\\Fabio\Desenvolvimento\\Varejo\\PPT\\TemplateR.pptx"),
    "V": Path(r"C:\\Fabio\Desenvolvimento\\Varejo\\PPT\\TemplateV.pptx"),
}

SHEET_NAME = "Dados"
COLUNA_FILTRO = "Nº Reguladora"

REQUIRED_COLUMNS = [
    "Nº Reguladora",
    "Prejuizo Apurado",
    "Valor do Embarque",
    "Data do Sinistro",
    "Causa Final",
    "Origem Ajustado",
    "Cidade - Destino",
    "UF - Origem",
    "UF - Destino",
    "Transportador",
    "Placa",
    "Motorista",
    "N_Carga",
    "QTDE_CARGAS_MOT",
    "Ação",
    # datas/etapas (podem vir vazias, mas checamos se existir)
    "ENCOSTA_EM_DOCA",
    "INICIO_CARREGAMENTO",
    "FIM_CARREGAMENTO",
    "EMISSAO_NF",
    "INICIO_VIAGEM",
    "CHEGADA_EM_LOJA",
]


# ============================================================
# LOGGER (console + queue para UI)
# ============================================================

class QueueHandler(logging.Handler):
    def __init__(self, q: Queue):
        super().__init__()
        self.q = q

    def emit(self, record: logging.LogRecord) -> None:
        try:
            msg = self.format(record)
            self.q.put(("log", msg))
        except Exception:
            pass


def build_logger(queue: Queue) -> logging.Logger:
    logger = logging.getLogger("sinistro_gui")
    logger.setLevel(logging.INFO)
    logger.handlers.clear()
    logger.propagate = False

    fmt = logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")

    # console
    ch = logging.StreamHandler(sys.stdout)
    ch.setLevel(logging.INFO)
    ch.setFormatter(fmt)
    logger.addHandler(ch)

    # queue -> UI
    qh = QueueHandler(queue)
    qh.setLevel(logging.INFO)
    qh.setFormatter(fmt)
    logger.addHandler(qh)

    return logger


# ============================================================
# FUNÇÕES DE NEGÓCIO (refatoradas)
# ============================================================

def format_datetime(value: Any) -> str:
    """
    Normaliza valores de data/hora e devolve string no formato "dd/mm/YYYY HH:MM".
    Aceita: datetime, pandas.Timestamp, string, número serial do Excel, NaN/None.
    """
    if value is None:
        return ""
    if value is None or pd.isna(value):  # Elimina o erro "NaTType does not support strftime" quando o campo estiver em branco
        return ""
    if isinstance(value, float) and pd.isna(value):
        return ""
    if isinstance(value, (datetime, pd.Timestamp)):
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y %H:%M")
    if isinstance(value, (int, float, np.floating)):
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y %H:%M")
        except Exception:
            return ""
    if isinstance(value, str):
        s = value.strip()
        if not s:
            return ""
        for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y %H:%M")
            except ValueError:
                continue
        return value
    return str(value)


def format_brl_currency(value: Any) -> str:
    """
    Retorna string no padrão BR: R$ 1.234,56
    Garante prefixo R$ uma única vez.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return "R$ 0,00"
    try:
        v = float(value)
        s = f"{v:,.2f}"  # 1,234.56
        s = s.replace(",", "X").replace(".", ",").replace("X", ".")  # 1.234,56
        return f"R$ {s}"
    except Exception:
        # se já veio texto, tenta normalizar prefixo
        txt = str(value).strip()
        if txt.upper().startswith("R$"):
            return txt
        return f"R$ {txt}" if txt else "R$ 0,00"


def safe_str(x: Any) -> str:
    if x is None:
        return ""
    if isinstance(x, float) and pd.isna(x):
        return ""
    return str(x).strip()


def load_excel(excel_path: Path, logger: logging.Logger) -> pd.DataFrame:
    if not excel_path.exists():
        raise FileNotFoundError(f"Excel não encontrado: {excel_path}")

    df = pd.read_excel(excel_path, sheet_name=SHEET_NAME)
    logger.info(f"Excel carregado: {excel_path}")
    logger.info(f"Linhas: {len(df)} | Colunas: {len(df.columns)}")

    # valida colunas essenciais (as que realmente usamos na geração)
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    if missing:
        # não necessariamente fatal para TODAS, mas é bom avisar com clareza
        # aqui vamos tratar como erro, por ser o comportamento mais seguro
        raise ValueError(f"Colunas ausentes na aba '{SHEET_NAME}': {missing}")

    return df


def list_sinistros(df: pd.DataFrame) -> List[str]:
    if COLUNA_FILTRO not in df.columns:
        return []
    vals = df[COLUNA_FILTRO].dropna().unique().tolist()
    # normaliza para string (mas mantendo número)
    out = []
    for v in vals:
        if isinstance(v, (int, np.integer)):
            out.append(str(int(v)))
        elif isinstance(v, float) and v.is_integer():
            out.append(str(int(v)))
        else:
            out.append(str(v))
    # remove duplicados mantendo ordem
    seen = set()
    uniq = []
    for x in out:
        if x not in seen:
            seen.add(x)
            uniq.append(x)
    return uniq


def get_sinistro_row(df: pd.DataFrame, numero_sinistro: int) -> pd.Series:
    # no excel pode vir float. Normalizamos comparando com int.
    col = df[COLUNA_FILTRO]
    mask = col.notna() & (col.astype(str).str.replace(".0", "", regex=False) == str(numero_sinistro))
    filtered = df[mask]
    if filtered.empty:
        # fallback: tentar comparação numérica direta
        try:
            filtered = df[df[COLUNA_FILTRO].notna() & (df[COLUNA_FILTRO].astype(float) == float(numero_sinistro))]
        except Exception:
            pass
    if filtered.empty:
        raise ValueError(f"Sinistro {numero_sinistro} não encontrado.")
    return filtered.iloc[0]


def extract_template_info(prs: Presentation) -> Dict[str, str]:
    """
    Tenta ler do template alguns textos para mapear substituições específicas.
    Mantém a sua lógica original (busca 'SINISTRO ...').
    """
    template_info: Dict[str, str] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                t = shape.text_frame.text or ""
                if "SINISTRO" in t:
                    m = re.search(r"SINISTRO (\d+\.?\d*)", t)
                    if m:
                        template_info["numero_sinistro"] = m.group(1)

                    m = re.search(r"– ([^–]+) –", t)
                    if m:
                        template_info["cidade_origem"] = m.group(1).strip()

                    m = re.search(r"\(([^)]+)\)", t)
                    if m:
                        template_info["destino_template"] = m.group(1).strip()
    return template_info


def replace_text_in_shape(shape, replacements: Dict[str, Any], titulo_slide: str) -> bool:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return False

    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            new_text = original_text

            for key, value in replacements.items():
                if key and key in new_text:
                    new_text = new_text.replace(key, str(value))
                    replaced = True

            if "SINISTRO" in new_text and "BACKOFFICE" in new_text:
                new_text = f"{titulo_slide}\nBACKOFFICE - SINISTROS"
                replaced = True

            if new_text != original_text:
                run.text = new_text

    return replaced


def criar_apresentacao_sinistro(
    df: pd.DataFrame,
    numero_sinistro: int,
    template_path: Path,
    output_dir: Path,
    logger: logging.Logger,
) -> Path:
    """
    Gera PPT aplicando substituições no template.
    """
    if not template_path.exists():
        raise FileNotFoundError(f"Template não encontrado: {template_path}")
    if not output_dir.exists():
        output_dir.mkdir(parents=True, exist_ok=True)

    sinistro_data = get_sinistro_row(df, numero_sinistro)

    prejuizo = sinistro_data.get("Prejuizo Apurado", None)
    valor_embarque = sinistro_data.get("Valor do Embarque", None)

    prejuizo_formatado = format_brl_currency(prejuizo)
    valor_embarque_formatado = format_brl_currency(valor_embarque)

    # Data sinistro dd/mm/yyyy (como seu original)
    data_sinistro = sinistro_data.get("Data do Sinistro", "")
    if isinstance(data_sinistro, (datetime, pd.Timestamp)):
        dt = data_sinistro.to_pydatetime() if isinstance(data_sinistro, pd.Timestamp) else data_sinistro
        data_formatada = dt.strftime("%d/%m/%Y")
    elif isinstance(data_sinistro, str):
        try:
            data_formatada = datetime.strptime(data_sinistro.strip(), "%d/%m/%Y").strftime("%d/%m/%Y")
        except Exception:
            data_formatada = safe_str(data_sinistro)
    else:
        data_formatada = safe_str(data_sinistro)

    # Carrega template
    prs = Presentation(str(template_path))
    logger.info(f"Template carregado: {template_path}")

    template_info = extract_template_info(prs)
    logger.info(f"Informações extraídas do template: {template_info}")

    causa = safe_str(sinistro_data.get("Causa Final", ""))
    cidade_origem = safe_str(sinistro_data.get("Origem Ajustado", ""))
    cidade_destino = safe_str(sinistro_data.get("Cidade - Destino", ""))
    uf_origem = safe_str(sinistro_data.get("UF - Origem", ""))
    uf_destino = safe_str(sinistro_data.get("UF - Destino", ""))

    titulo_slide = f"SINISTRO {numero_sinistro} – {causa} – {cidade_origem} – ({cidade_destino})"

    # Mapeamento dados
    sinistro_info: Dict[str, Any] = {
        template_info.get("numero_sinistro", "65.329"): str(numero_sinistro),
        template_info.get("cidade_origem", "ITAPECERICA DA SERRA"): cidade_origem,

        "SINISTRO ": f"SINISTRO {numero_sinistro}",
        "CAUSA": causa,
        "Info_CidadeOrigem": cidade_origem,
        "Info_Cid_Origem": safe_str(sinistro_data.get("Cidade Origem", "")),
        "Info_Uf_Origem": uf_origem,
        "Info_Destino": cidade_destino,
        "Info_Uf_Destino": uf_destino,
        "Info_Transp": safe_str(sinistro_data.get("Transportador", "")),
        "Info_Placa": safe_str(sinistro_data.get("Placa", "")),
        "Info_mot": safe_str(sinistro_data.get("Motorista", "")),
        # aqui NÃO duplica "R$"
        "Info_VlCarga": valor_embarque_formatado,
        "Info_Prejuizo": prejuizo_formatado,
        "Info_DT_Sinistro": data_formatada,
        "Info_Carga": safe_str(sinistro_data.get("N_Carga", "")),
        "Info_Qte_Cargas": safe_str(sinistro_data.get("QTDE_CARGAS_MOT", "")),
        "Info_Desc": safe_str(sinistro_data.get("Ação", "")),

        "Info_Doca":           format_datetime(sinistro_data.get("ENCOSTA_EM_DOCA", "")),
        "Info_Inicio_Carreg":  format_datetime(sinistro_data.get("INICIO_CARREGAMENTO", "")),
        "Info_Fim_Carreg":     format_datetime(sinistro_data.get("FIM_CARREGAMENTO", "")),
        "Info_Emissao_NF":     format_datetime(sinistro_data.get("EMISSAO_NF", "")),
        "Info_Inicio_Viagem":  format_datetime(sinistro_data.get("INICIO_VIAGEM", "")),
        "Info_Chegada":        format_datetime(sinistro_data.get("CHEGADA_EM_LOJA", "")),
    }

    # Substituição da descrição em tabelas (mantendo intenção original)
    descricao = safe_str(sinistro_data.get("Ação", ""))
    if descricao:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            txt = cell.text_frame.text if cell.text_frame is not None else ""
                            if txt and "Veículo saiu carregado" in txt:
                                sinistro_info[txt] = descricao
                            elif txt and "Resumo Ocorrência" in txt:
                                # tenta trocar célula adjacente (heurística)
                                for adj_cell in row.cells:
                                    if adj_cell is not cell:
                                        adj_txt = adj_cell.text_frame.text if adj_cell.text_frame else ""
                                        if adj_txt:
                                            sinistro_info[adj_txt] = descricao

    logger.info(f"Mapeamento criado com {len(sinistro_info)} itens")

    total_replacements = 0
    for i, slide in enumerate(prs.slides, start=1):
        slide_repl = 0
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                if replace_text_in_shape(shape, sinistro_info, titulo_slide):
                    slide_repl += 1

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_text_in_shape(cell, sinistro_info, titulo_slide):
                            slide_repl += 1

        logger.info(f"Slide {i}: {slide_repl} substituições")
        total_replacements += slide_repl

    logger.info(f"Total substituições: {total_replacements}")

    output_path = output_dir / f"Sinistro_{numero_sinistro}_Final_v2.pptx"
    prs.save(str(output_path))
    logger.info(f"Apresentação salva: {output_path}")
    return output_path


# ============================================================
# UI (Tkinter futurista escuro)
# ============================================================

@dataclass
class UiState:
    df: Optional[pd.DataFrame] = None
    last_output_ppt: Optional[Path] = None


class SinistroGUI:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("Sinistros | Gerador PPT (Tkinter)")

        self.queue: Queue = Queue()
        self.logger = build_logger(self.queue)
        self.state = UiState()

        self.colors = {
            "bg": "#0b0f1a",
            "panel": "#0f172a",
            "entry": "#111827",
            "fg": "#e6e6e6",
            "muted": "#94a3b8",
            "accent": "#00e5ff",
            "highlight": "#7c3aed",
            "danger": "#ef4444",
        }

        self._setup_style()
        self._build_layout()
        self._poll_queue()

    def _setup_style(self):
        self.root.configure(bg=self.colors["bg"])
        style = ttk.Style(self.root)
        style.theme_use("clam")

        style.configure(".", background=self.colors["bg"], foreground=self.colors["fg"])

        style.configure("TFrame", background=self.colors["bg"])
        style.configure("Panel.TFrame", background=self.colors["panel"])
        style.configure("TLabel", background=self.colors["bg"], foreground=self.colors["fg"])
        style.configure("Title.TLabel", font=("Segoe UI", 14, "bold"), foreground=self.colors["accent"])
        style.configure("Muted.TLabel", foreground=self.colors["muted"])

        style.configure(
            "TButton",
            background=self.colors["entry"],
            foreground=self.colors["fg"],
            bordercolor=self.colors["accent"],
            lightcolor=self.colors["accent"],
            darkcolor=self.colors["accent"],
            focuscolor=self.colors["highlight"],
            padding=(10, 6),
        )
        style.map("TButton", background=[("active", self.colors["highlight"])])

        style.configure(
            "TEntry",
            fieldbackground=self.colors["entry"],
            foreground=self.colors["fg"],
            insertcolor=self.colors["accent"],
        )

        style.configure("TRadiobutton", background=self.colors["bg"], foreground=self.colors["fg"])
        style.configure("TCheckbutton", background=self.colors["bg"], foreground=self.colors["fg"])

        style.configure(
            "Treeview",
            background=self.colors["panel"],
            fieldbackground=self.colors["panel"],
            foreground=self.colors["fg"],
            bordercolor=self.colors["accent"],
            rowheight=24,
        )
        style.map("Treeview", background=[("selected", "#1f2937")], foreground=[("selected", self.colors["accent"])])

        style.configure("Horizontal.TProgressbar", background=self.colors["accent"], troughcolor=self.colors["entry"])

        style.configure("TNotebook", background=self.colors["bg"], bordercolor=self.colors["accent"])
        style.configure("TNotebook.Tab", background=self.colors["entry"], foreground=self.colors["fg"], padding=(12, 6))
        style.map("TNotebook.Tab", background=[("selected", self.colors["highlight"])])

    def _build_layout(self):
        self.root.geometry("1100x700")
        self.root.minsize(1000, 650)

        # Header
        header = ttk.Frame(self.root, style="Panel.TFrame")
        header.pack(fill="x", padx=12, pady=(12, 8))

        ttk.Label(header, text="GERADOR DE PPT - SINISTROS - Fábio A Cordeiro", style="Title.TLabel").pack(anchor="w", padx=12, pady=(10, 2))
        ttk.Label(
            header,
            text="Interface interativa | Tema escuro | Execução em thread | Logs em tempo real",
            style="Muted.TLabel",
        ).pack(anchor="w", padx=12, pady=(0, 10))

        # Body split (left controls / right preview)
        body = ttk.Frame(self.root)
        body.pack(fill="both", expand=True, padx=12, pady=(0, 12))

        left = ttk.Frame(body, style="Panel.TFrame")
        right = ttk.Frame(body, style="Panel.TFrame")
        left.pack(side="left", fill="y", padx=(0, 10), pady=0)
        right.pack(side="right", fill="both", expand=True, padx=0, pady=0)

        # LEFT - Controls notebook
        nb = ttk.Notebook(left)
        nb.pack(fill="both", expand=True, padx=10, pady=10)

        tab_files = ttk.Frame(nb)
        tab_filter = ttk.Frame(nb)
        tab_actions = ttk.Frame(nb)

        nb.add(tab_files, text="Arquivos")
        nb.add(tab_filter, text="Filtro")
        nb.add(tab_actions, text="Ações")

        # RIGHT - Preview + log
        preview_title = ttk.Frame(right)
        preview_title.pack(fill="x", padx=10, pady=(10, 6))
        ttk.Label(preview_title, text="Prévia do Sinistro", style="Title.TLabel").pack(anchor="w")

        self.tree = ttk.Treeview(
            right,
            columns=("campo", "valor"),
            show="headings",
            height=10,
        )
        self.tree.heading("campo", text="Campo")
        self.tree.heading("valor", text="Valor")
        self.tree.column("campo", width=260, anchor="w")
        self.tree.column("valor", width=560, anchor="w")
        self.tree.pack(fill="x", padx=10, pady=(0, 10))

        log_title = ttk.Frame(right)
        log_title.pack(fill="x", padx=10, pady=(0, 6))
        ttk.Label(log_title, text="Logs", style="Title.TLabel").pack(anchor="w")

        self.log_text = scrolledtext.ScrolledText(
            right,
            height=14,
            bg=self.colors["entry"],
            fg=self.colors["fg"],
            insertbackground=self.colors["accent"],
            relief="flat",
            highlightthickness=1,
            highlightbackground=self.colors["accent"],
        )
        self.log_text.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        # Status bar
        self.status_var = tk.StringVar(value="Pronto.")
        status = ttk.Frame(self.root, style="Panel.TFrame")
        status.pack(fill="x", padx=12, pady=(0, 12))
        ttk.Label(status, textvariable=self.status_var).pack(anchor="w", padx=12, pady=8)

        # ----------------------------
        # Tab: Arquivos
        # ----------------------------
        self.excel_var = tk.StringVar(value=str(DEFAULT_EXCEL))
        self.output_var = tk.StringVar(value=str(DEFAULT_OUTPUT_DIR))
        self.tp_var = tk.StringVar(value="A")
        self.manual_template_var = tk.BooleanVar(value=False)
        self.manual_template_path_var = tk.StringVar(value="")

        row = 0
        ttk.Label(tab_files, text="Excel de origem (aba 'Dados')").grid(row=row, column=0, sticky="w", padx=10, pady=(10, 4))
        row += 1

        excel_row = ttk.Frame(tab_files)
        excel_row.grid(row=row, column=0, sticky="ew", padx=10, pady=(0, 10))
        excel_row.columnconfigure(0, weight=1)
        ttk.Entry(excel_row, textvariable=self.excel_var, width=56).grid(row=0, column=0, sticky="ew")
        ttk.Button(excel_row, text="Procurar", command=self._pick_excel).grid(row=0, column=1, padx=(8, 0))
        row += 1

        ttk.Label(tab_files, text="Tipo de sinistro (template A/R/V)").grid(row=row, column=0, sticky="w", padx=10, pady=(0, 4))
        row += 1

        tp_row = ttk.Frame(tab_files)
        tp_row.grid(row=row, column=0, sticky="w", padx=10, pady=(0, 10))
        ttk.Radiobutton(tp_row, text="A - Acidente/Avaria/Tombamento/Colisão", variable=self.tp_var, value="A").grid(row=0, column=0, sticky="w")
        ttk.Radiobutton(tp_row, text="R - Roubo/Furto", variable=self.tp_var, value="R").grid(row=1, column=0, sticky="w")
        ttk.Radiobutton(tp_row, text="V - Variação Temperatura", variable=self.tp_var, value="V").grid(row=2, column=0, sticky="w")
        row += 1

        manual_row = ttk.Frame(tab_files)
        manual_row.grid(row=row, column=0, sticky="ew", padx=10, pady=(0, 10))
        manual_row.columnconfigure(1, weight=1)
        ttk.Checkbutton(
            manual_row,
            text="Escolher template manualmente",
            variable=self.manual_template_var,
            command=self._toggle_manual_template,
        ).grid(row=0, column=0, sticky="w")
        self.manual_entry = ttk.Entry(manual_row, textvariable=self.manual_template_path_var, state="disabled")
        self.manual_entry.grid(row=0, column=1, sticky="ew", padx=(8, 8))
        self.btn_pick_template = ttk.Button(manual_row, text="Procurar", command=self._pick_template, state="disabled")
        self.btn_pick_template.grid(row=0, column=2)
        row += 1

        ttk.Label(tab_files, text="Pasta de saída").grid(row=row, column=0, sticky="w", padx=10, pady=(0, 4))
        row += 1

        out_row = ttk.Frame(tab_files)
        out_row.grid(row=row, column=0, sticky="ew", padx=10, pady=(0, 10))
        out_row.columnconfigure(0, weight=1)
        ttk.Entry(out_row, textvariable=self.output_var).grid(row=0, column=0, sticky="ew")
        ttk.Button(out_row, text="Procurar", command=self._pick_output_dir).grid(row=0, column=1, padx=(8, 0))
        row += 1

        # ----------------------------
        # Tab: Filtro
        # ----------------------------
        self.sinistro_var = tk.StringVar(value="")
        self.sinistros_combo_var = tk.StringVar(value="")
        self._vcmd_num = (self.root.register(self._validate_only_digits), "%P")

        ttk.Label(tab_filter, text="Número do sinistro (Nº Reguladora)").grid(row=0, column=0, sticky="w", padx=10, pady=(10, 4))
        sin_row = ttk.Frame(tab_filter)
        sin_row.grid(row=1, column=0, sticky="ew", padx=10, pady=(0, 10))
        sin_row.columnconfigure(0, weight=1)
        ttk.Entry(sin_row, textvariable=self.sinistro_var, validate="key", validatecommand=self._vcmd_num).grid(row=0, column=0, sticky="ew")
        ttk.Button(sin_row, text="Carregar/Validar Excel", command=self._load_validate_excel).grid(row=0, column=1, padx=(8, 0))

        ttk.Label(tab_filter, text="Sinistros disponíveis (após carregar)").grid(row=2, column=0, sticky="w", padx=10, pady=(0, 4))
        combo_row = ttk.Frame(tab_filter)
        combo_row.grid(row=3, column=0, sticky="ew", padx=10, pady=(0, 10))
        combo_row.columnconfigure(0, weight=1)
        self.sinistro_combo = ttk.Combobox(combo_row, textvariable=self.sinistros_combo_var, state="disabled")
        self.sinistro_combo.grid(row=0, column=0, sticky="ew")
        self.sinistro_combo.bind("<<ComboboxSelected>>", self._on_combo_selected)
        ttk.Button(combo_row, text="Pré-visualizar", command=self._preview_sinistro).grid(row=0, column=1, padx=(8, 0))

        ttk.Label(tab_filter, text="Dica: você pode digitar o sinistro ou escolher no dropdown.").grid(
            row=4, column=0, sticky="w", padx=10, pady=(0, 10)
        )

        # ----------------------------
        # Tab: Ações
        # ----------------------------
        ttk.Label(tab_actions, text="Ações").grid(row=0, column=0, sticky="w", padx=10, pady=(10, 6))

        btn_row = ttk.Frame(tab_actions)
        btn_row.grid(row=1, column=0, sticky="ew", padx=10, pady=(0, 10))
        ttk.Button(btn_row, text="Gerar PPT", command=self._generate_ppt_clicked).grid(row=0, column=0, padx=(0, 8))
        ttk.Button(btn_row, text="Abrir pasta de saída", command=self._open_output_dir).grid(row=0, column=1)

        self.progress = ttk.Progressbar(tab_actions, mode="indeterminate", style="Horizontal.TProgressbar")
        self.progress.grid(row=2, column=0, sticky="ew", padx=10, pady=(0, 10))

        ttk.Label(
            tab_actions,
            text="A geração roda em thread (não trava a UI). Erros aparecem em popup + log.",
        ).grid(row=3, column=0, sticky="w", padx=10, pady=(0, 10))

        # Expand behavior for tabs
        for t in (tab_files, tab_filter, tab_actions):
            t.columnconfigure(0, weight=1)

    # ----------------------------
    # UI helpers
    # ----------------------------
    def _set_status(self, msg: str):
        self.status_var.set(msg)

    def _append_log(self, msg: str):
        self.log_text.insert("end", msg + "\n")
        self.log_text.see("end")

    def _poll_queue(self):
        try:
            while True:
                kind, payload = self.queue.get_nowait()
                if kind == "log":
                    self._append_log(payload)
                elif kind == "done":
                    # payload: dict com status
                    self._on_generation_done(payload)
        except Empty:
            pass
        self.root.after(120, self._poll_queue)

    def _validate_only_digits(self, new_value: str) -> bool:
        return new_value.isdigit() or new_value == ""

    def _pick_excel(self):
        path = filedialog.askopenfilename(
            title="Selecione o Excel",
            filetypes=[("Excel", "*.xlsx *.xls")],
        )
        if path:
            self.excel_var.set(path)

    def _pick_output_dir(self):
        path = filedialog.askdirectory(title="Selecione a pasta de saída")
        if path:
            self.output_var.set(path)

    def _toggle_manual_template(self):
        is_manual = self.manual_template_var.get()
        self.manual_entry.configure(state="normal" if is_manual else "disabled")
        self.btn_pick_template.configure(state="normal" if is_manual else "disabled")

    def _pick_template(self):
        path = filedialog.askopenfilename(
            title="Selecione o template PPTX",
            filetypes=[("PowerPoint", "*.pptx")],
        )
        if path:
            self.manual_template_path_var.set(path)

    def _resolve_template_path(self) -> Path:
        if self.manual_template_var.get():
            p = Path(self.manual_template_path_var.get().strip())
            if not p.exists():
                raise FileNotFoundError(f"Template manual não encontrado: {p}")
            return p

        tp = (self.tp_var.get() or "A").strip().upper()[:1]
        if tp not in TEMPLATE_PATHS:
            tp = "A"
        p = TEMPLATE_PATHS[tp]
        if not p.exists():
            raise FileNotFoundError(f"Template padrão não encontrado para '{tp}': {p}")
        return p

    def _load_validate_excel(self):
        excel_path = Path(self.excel_var.get().strip())
        try:
            self._set_status("Carregando Excel...")
            self.logger.info("Iniciando carregamento do Excel...")
            df = load_excel(excel_path, self.logger)

            self.state.df = df
            sins = list_sinistros(df)
            self.sinistro_combo.configure(state="readonly")
            self.sinistro_combo["values"] = sins
            self._set_status(f"Excel OK | {len(df)} linhas | {len(sins)} sinistros disponíveis")
            self.logger.info(f"Sinistros disponíveis: {len(sins)}")

        except Exception as e:
            self._set_status("Erro ao carregar Excel.")
            messagebox.showerror("Erro", str(e))
            self.logger.error(str(e))

    def _on_combo_selected(self, _event=None):
        val = self.sinistros_combo_var.get()
        if val:
            self.sinistro_var.set(val)

    def _preview_sinistro(self):
        if self.state.df is None:
            messagebox.showwarning("Aviso", "Carregue/valide o Excel primeiro.")
            return

        n_txt = self.sinistro_var.get().strip()
        if not n_txt:
            messagebox.showwarning("Aviso", "Informe o Nº Reguladora.")
            return

        try:
            numero = int(n_txt)
            row = get_sinistro_row(self.state.df, numero)

            # limpa tree
            for iid in self.tree.get_children():
                self.tree.delete(iid)

            # campos principais (como você pediu)
            campos = [
                "Nº Reguladora",
                "Data do Sinistro",
                "Transportador",
                "Placa",
                "Prejuizo Apurado",
                "Valor do Embarque",
                "Causa Final",
                "Origem Ajustado",
                "Cidade - Destino",
            ]

            for c in campos:
                v = row.get(c, "")
                if c in ("Prejuizo Apurado", "Valor do Embarque"):
                    v = format_brl_currency(v)
                elif c == "Data do Sinistro":
                    # dd/mm/yyyy se for datetime
                    if isinstance(v, (datetime, pd.Timestamp)):
                        dt = v.to_pydatetime() if isinstance(v, pd.Timestamp) else v
                        v = dt.strftime("%d/%m/%Y")
                    v = safe_str(v)
                else:
                    v = safe_str(v)

                self.tree.insert("", "end", values=(c, v))

            self._set_status(f"Prévia carregada: sinistro {numero}")
            self.logger.info(f"Prévia carregada para sinistro {numero}")

        except Exception as e:
            messagebox.showerror("Erro", str(e))
            self.logger.error(str(e))
            self._set_status("Erro na prévia.")

    def _generate_ppt_clicked(self):
        if self.state.df is None:
            messagebox.showwarning("Aviso", "Carregue/valide o Excel primeiro.")
            return

        n_txt = self.sinistro_var.get().strip()
        if not n_txt:
            messagebox.showwarning("Aviso", "Informe o Nº Reguladora.")
            return

        try:
            numero = int(n_txt)
        except ValueError:
            messagebox.showerror("Erro", "O Nº Reguladora deve ser numérico.")
            return

        try:
            template_path = self._resolve_template_path()
            output_dir = Path(self.output_var.get().strip())
        except Exception as e:
            messagebox.showerror("Erro", str(e))
            self.logger.error(str(e))
            return

        # inicia progresso
        self.progress.start(12)
        self._set_status("Gerando PPT...")

        def worker():
            try:
                self.logger.info(f"Tipo de sinistro informado: {self.tp_var.get()}")
                out = criar_apresentacao_sinistro(
                    df=self.state.df,
                    numero_sinistro=numero,
                    template_path=template_path,
                    output_dir=output_dir,
                    logger=self.logger,
                )
                self.queue.put(("done", {"ok": True, "path": str(out)}))
            except Exception as e:
                self.queue.put(("done", {"ok": False, "error": str(e)}))

        threading.Thread(target=worker, daemon=True).start()

    def _on_generation_done(self, payload: dict):
        self.progress.stop()
        if payload.get("ok"):
            p = Path(payload["path"])
            self.state.last_output_ppt = p
            self._set_status(f"PPT gerado: {p.name}")
            messagebox.showinfo("OK", f"PPT gerado com sucesso:\n{p}")
        else:
            err = payload.get("error", "Erro desconhecido")
            self._set_status("Erro na geração do PPT.")
            messagebox.showerror("Erro", err)
            self.logger.error(err)

    def _open_output_dir(self):
        out_dir = Path(self.output_var.get().strip())
        if out_dir.exists():
            try:
                os.startfile(str(out_dir))
            except Exception as e:
                messagebox.showerror("Erro", str(e))
        else:
            messagebox.showerror("Erro", "Pasta de saída não existe.")


# ============================================================
# MAIN
# ============================================================

def main():
    root = tk.Tk()
    app = SinistroGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()