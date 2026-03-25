# pip install python-pptx pandas openpyxl numpy
# Criado por Fábio A Cordeiro em 24/08/2025
# UI Tkinter + correção definitiva do erro: NaTType does not support strftime

import os
import re
import sys
import queue
import threading
from datetime import datetime, timedelta

import numpy as np
import pandas as pd
from pptx import Presentation

# -----------------------------
# Defaults (editáveis na UI)
# -----------------------------
DEFAULT_EXCEL = r"C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\SINISTROS1.xlsx"
DEFAULT_TEMPLATE_A = r"C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\TemplateA.pptx"
DEFAULT_TEMPLATE_R = r"C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\TemplateR.pptx"
DEFAULT_TEMPLATE_V = r"C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\TemplateV.pptx"


# -----------------------------
# Utilitários (robustos)
# -----------------------------
def safe_filename(value: object) -> str:
    s = str(value) if value is not None else ""
    s = s.strip()
    s = re.sub(r"[<>:\"/\|?*\n\r\t]", "_", s)
    s = re.sub(r"\s+", " ", s)
    return s or "SemNome"


def format_brl_currency(value) -> str:
    # Retorna SEMPRE string; se numérico -> "R$ 1.234,56"
    if value is None or pd.isna(value):
        return ""
    if isinstance(value, (int, float, np.floating)) and pd.notna(value):
        return f"R$ {float(value):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    s = str(value).strip()
    if s.lower() in ("nan", "nat", "n/a", "none"):
        return ""
    return s


def format_datetime(value) -> str:
    """
    Normaliza valores de data/hora e devolve string no formato "dd/mm/YYYY HH:MM".
    Nunca chama strftime() em NaT.
    """
    if value is None or pd.isna(value):
        return ""

    # datetime / Timestamp
    if isinstance(value, (datetime, pd.Timestamp)):
        if pd.isna(value):
            return ""
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y %H:%M")

    # número serial do Excel
    if isinstance(value, (int, float, np.floating)):
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y %H:%M")
        except Exception:
            return ""

    # string
    if isinstance(value, str):
        s = value.strip()
        if not s or s.lower() in ("nan", "nat", "n/a", "none"):
            return ""
        for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y %H:%M")
            except ValueError:
                pass
        return s

    return str(value)


def format_date_ddmmyyyy(value) -> str:
    """
    Formata uma data para dd/mm/yyyy sem quebrar com NaT.
    """
    if value is None or pd.isna(value):
        return ""
    if isinstance(value, (datetime, pd.Timestamp)):
        if pd.isna(value):
            return ""
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y")
    if isinstance(value, (int, float, np.floating)):
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y")
        except Exception:
            return ""
    if isinstance(value, str):
        s = value.strip()
        if not s or s.lower() in ("nan", "nat", "n/a", "none"):
            return ""
        for fmt in ("%d/%m/%Y", "%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y")
            except ValueError:
                pass
        return s
    return str(value)


def replace_text_in_shape(shape, replacements: dict, titulo_slide: str) -> bool:
    """
    Substitui texto mantendo formatação (substituição por run).
    """
    if not hasattr(shape, "text_frame") or not shape.has_text_frame:
        return False

    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            new_text = original_text

            for key, value in replacements.items():
                if key and key in new_text:
                    new_text = new_text.replace(key, str(value))
                    replaced = True

            # Substituir título completo do slide (regra do script original)
            if "SINISTRO" in new_text and "BACKOFFICE" in new_text:
                new_text = f"{titulo_slide}\nBACKOFFICE - SINISTROS"
                replaced = True

            if new_text != original_text:
                run.text = new_text

    return replaced


# -----------------------------
# Núcleo: criação do PPT
# -----------------------------
def criar_apresentacao_sinistro(
    numero_sinistro: int | None,
    excel_source_path: str,
    template_path: str,
    output_dir: str,
    log_func=None,
) -> str | None:
    def log(msg: str):
        if log_func:
            log_func(msg)

    # Ler o Excel (aba Dados, como no seu script original)
    try:
        df = pd.read_excel(excel_source_path, sheet_name="Dados")
        log(f"Excel carregado: {excel_source_path} | linhas: {len(df)}")
    except Exception as e:
        log(f"Erro ao ler o Excel: {e}")
        return None

    if df.empty:
        log("A planilha (aba 'Dados') está vazia.")
        return None

    # Conversão preventiva de colunas de datas (evita mistura de tipos)
    cols_dt = [
        "Data do Sinistro",
        "ENCOSTA_EM_DOCA",
        "INICIO_CARREGAMENTO",
        "FIM_CARREGAMENTO",
        "EMISSAO_NF",
        "INICIO_VIAGEM",
        "CHEGADA_EM_LOJA",
    ]
    for c in cols_dt:
        if c in df.columns:
            df[c] = pd.to_datetime(df[c], errors="coerce")

    # Filtrar somente linhas com Nº Reguladora válido (como original)
    coluna_filtro = "Nº Reguladora"
    if coluna_filtro in df.columns:
        df_filtrado = df[df[coluna_filtro].notna()]
    else:
        log(f"Coluna obrigatória não encontrada no Excel: {coluna_filtro}")
        return None

    # Se numero_sinistro foi fornecido, filtra
    if numero_sinistro is not None:
        sinistro_df = df_filtrado[df_filtrado["Nº Reguladora"] == numero_sinistro]
        if sinistro_df.empty:
            log(f"Sinistro {numero_sinistro} não encontrado na aba 'Dados'.")
            return None
    else:
        # fallback: primeira linha filtrada
        sinistro_df = df_filtrado.copy()
        if sinistro_df.empty:
            log("Nenhum registro com Nº Reguladora válido encontrado.")
            return None

    # Salva Excel filtrado (como seu script, mas no output_dir)
    excel_filtrado_path = os.path.join(output_dir, "SINISTROS_FILTRADO.xlsx")
    try:
        sinistro_df.to_excel(excel_filtrado_path, index=False)
        log(f"Excel filtrado gerado: {excel_filtrado_path}")
    except Exception as e:
        log(f"Erro ao salvar Excel filtrado: {e}")
        return None

    # Usa primeira linha do sinistro filtrado
    sinistro_data = sinistro_df.iloc[0]

    # Define numero_sinistro se não veio
    if numero_sinistro is None:
        try:
            numero_sinistro = int(sinistro_data.get("Nº Reguladora"))
        except Exception:
            numero_sinistro = None

    numero_sinistro_str = safe_filename(numero_sinistro if numero_sinistro is not None else "Desconhecido")
    log(f"Processando sinistro: {numero_sinistro_str}")

    # Formata valores
    prejuizo_formatado = format_brl_currency(sinistro_data.get("Prejuizo Apurado", ""))
    valor_embarque_formatado = format_brl_currency(sinistro_data.get("Valor do Embarque", ""))

    # Formata data do sinistro (NUNCA quebra com NaT)
    data_formatada = format_date_ddmmyyyy(sinistro_data.get("Data do Sinistro", None))

    # Carrega template PPT
    try:
        prs = Presentation(template_path)
        log(f"Template carregado: {template_path}")
    except Exception as e:
        log(f"Erro ao carregar o template PPT: {e}")
        return None

    # Extrair informações do template para substituição (como original)
    template_info = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                txt = shape.text_frame.text or ""
                if "SINISTRO" in txt:
                    match = re.search(r"SINISTRO (\d+\.?\d*)", txt)
                    if match:
                        template_info["numero_sinistro"] = match.group(1)

                    match = re.search(r"– ([^–]+) –", txt)
                    if match:
                        template_info["cidade_origem"] = match.group(1).strip()

                    match = re.search(r"\(([^)]+)\)", txt)
                    if match:
                        template_info["destino_template"] = match.group(1).strip()

    log(f"Info extraída do template: {template_info}")

    # Cria título do slide
    causa = str(sinistro_data.get("Causa Final", "")).strip()
    cidade_origem = str(sinistro_data.get("Origem Ajustado", "")).strip()
    cidade_destino = str(sinistro_data.get("Cidade - Destino", "")).strip()
    uf_origem = str(sinistro_data.get("UF - Origem", "")).strip()
    uf_destino = str(sinistro_data.get("UF - Destino", "")).strip()

    titulo_slide = f"SINISTRO {numero_sinistro_str} – {causa} – {cidade_origem} – ({cidade_destino})"

    # Mapeamento de substituições (ajustes: sem duplicar "R$")
    sinistro_info = {
        # Substituições específicas do template
        template_info.get("numero_sinistro", "65.329"): str(numero_sinistro_str),
        template_info.get("cidade_origem", "ITAPECERICA DA SERRA"): cidade_origem,
        template_info.get("Complemento_info", " - "): cidade_destino,

        # Substituições gerais
        "SINISTRO ": f"SINISTRO {numero_sinistro_str}",
        "CAUSA": causa,
        "Info_CidadeOrigem": cidade_origem,
        "Info_Cid_Origem": str(sinistro_data.get("Cidade Origem", "")),
        "Info_Uf_Origem": uf_origem,
        "Info_Destino": cidade_destino,
        "Info_Cid_Destino": str(sinistro_data.get("Complemento_info", "")),
        "Info_Transp": str(sinistro_data.get("Transportador", "")),
        "Info_Placa": str(sinistro_data.get("Placa", "")),
        "Info_mot": str(sinistro_data.get("Motorista", "")),
        "Info_VlCarga": valor_embarque_formatado,      # já vem "R$ ..." se numérico
        "Info_Prejuizo": prejuizo_formatado,           # já vem "R$ ..." se numérico
        "Info_DT_Sinistro": data_formatada,
        "Info_Carga": str(sinistro_data.get("N_Carga", "")),
        "Info_Qte_Cargas": str(sinistro_data.get("QTDE_CARGAS_MOT", "")),
        "Info_Desc": str(sinistro_data.get("Ação", "")),

        # DATAS/HORAS formatadas (robustas)
        "Info_Doca": format_datetime(sinistro_data.get("ENCOSTA_EM_DOCA", "")),
        "Info_Inicio_Carreg": format_datetime(sinistro_data.get("INICIO_CARREGAMENTO", "")),
        "Info_Fim_Carreg": format_datetime(sinistro_data.get("FIM_CARREGAMENTO", "")),
        "Info_Emissao_NF": format_datetime(sinistro_data.get("EMISSAO_NF", "")),
        "Info_Inicio_Viagem": format_datetime(sinistro_data.get("INICIO_VIAGEM", "")),
        "Info_Chegada": format_datetime(sinistro_data.get("CHEGADA_EM_LOJA", "")),
    }

    # Descrição (Ação) substituindo texto em tabela, como seu original
    descricao = str(sinistro_data.get("Ação", "") or "")
    if descricao:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_txt = cell.text_frame.text if hasattr(cell, "text_frame") else ""
                            if cell_txt and "Veículo saiu carregado" in cell_txt:
                                sinistro_info[cell_txt] = descricao
                            elif cell_txt and "Resumo Ocorrência" in cell_txt:
                                if len(row.cells) > 1:
                                    for adj_cell in row.cells:
                                        if adj_cell != cell and hasattr(adj_cell, "text_frame") and adj_cell.text_frame.text:
                                            sinistro_info[adj_cell.text_frame.text] = descricao

    log(f"Mapeamento criado com {len(sinistro_info)} itens")

    # Substituição em slides
    total_replacements = 0
    for slide_idx, slide in enumerate(prs.slides):
        slide_replacements = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                if replace_text_in_shape(shape, sinistro_info, titulo_slide):
                    slide_replacements += 1

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_text_in_shape(cell, sinistro_info, titulo_slide):
                            slide_replacements += 1

        total_replacements += slide_replacements
        log(f"Slide {slide_idx + 1}: {slide_replacements} substituições")

    log(f"Total de substituições: {total_replacements}")

    # Salva apresentação
    output_path = os.path.join(output_dir, f"Sinistro_{numero_sinistro_str}_Final_v2.pptx")
    try:
        prs.save(output_path)
        log(f"PPT gerado com sucesso: {output_path}")
        return output_path
    except Exception as e:
        log(f"Erro ao salvar PPT: {e}")
        return None


# -----------------------------
# Tkinter UI (tema escuro/neon)
# -----------------------------
def run_ui():
    import tkinter as tk
    from tkinter import filedialog, ttk

    BG = "#0b0f1a"
    PANEL = "#0f172a"
    ENTRY_BG = "#111827"
    FG = "#dbeafe"
    MUTED = "#93c5fd"
    NEON = "#00e5ff"
    PURPLE = "#7c4dff"
    BTN_BG = "#121a2b"
    BTN_BG_HOVER = "#16213a"

    q = queue.Queue()

    root = tk.Tk()
    root.title("Gerador de PPT - Sinistros (Futurista)")
    root.geometry("980x680")
    root.configure(bg=BG)
    root.minsize(920, 620)

    style = ttk.Style()
    style.theme_use("clam")

    # Estilos ttk
    style.configure("TFrame", background=BG)
    style.configure("Card.TFrame", background=PANEL)
    style.configure("TLabel", background=PANEL, foreground=FG, font=("Segoe UI", 10))
    style.configure("Title.TLabel", background=BG, foreground=NEON, font=("Segoe UI Semibold", 16))
    style.configure("Hint.TLabel", background=PANEL, foreground=MUTED, font=("Segoe UI", 9))
    style.configure("TRadiobutton", background=PANEL, foreground=FG, font=("Segoe UI", 10))
    style.map("TRadiobutton", foreground=[("active", NEON)])

    style.configure("Neon.TButton", background=BTN_BG, foreground=FG, borderwidth=0, focusthickness=0,
                    font=("Segoe UI Semibold", 10), padding=(12, 8))
    style.map("Neon.TButton",
              background=[("active", BTN_BG_HOVER), ("pressed", "#0b1220")],
              foreground=[("active", NEON)])

    style.configure("Accent.TButton", background=PURPLE, foreground="white", borderwidth=0, focusthickness=0,
                    font=("Segoe UI Semibold", 10), padding=(12, 10))
    style.map("Accent.TButton",
              background=[("active", "#6d28d9"), ("pressed", "#4c1d95")])

    style.configure("TEntry", fieldbackground=ENTRY_BG, foreground=FG, insertcolor=FG, bordercolor="#1f2a44")
    style.configure("TCombobox", fieldbackground=ENTRY_BG, foreground=FG)

    # Variáveis
    excel_var = tk.StringVar(value=DEFAULT_EXCEL)
    sinistro_var = tk.StringVar(value="")
    tipo_var = tk.StringVar(value="A")

    template_a_var = tk.StringVar(value=DEFAULT_TEMPLATE_A)
    template_r_var = tk.StringVar(value=DEFAULT_TEMPLATE_R)
    template_v_var = tk.StringVar(value=DEFAULT_TEMPLATE_V)

    outdir_var = tk.StringVar(value="")  # default: pasta do excel

    # Layout
    header = ttk.Frame(root)
    header.pack(fill="x", padx=18, pady=(14, 10))
    ttk.Label(header, text="Gerador de PPT de Sinistros - Created by: Fabio A Cordeiro", style="Title.TLabel").pack(anchor="w")

    content = ttk.Frame(root)
    content.pack(fill="both", expand=True, padx=18, pady=(0, 18))
    content.columnconfigure(0, weight=3)
    content.columnconfigure(1, weight=2)
    content.rowconfigure(0, weight=1)

    left = ttk.Frame(content, style="Card.TFrame")
    right = ttk.Frame(content, style="Card.TFrame")
    left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))
    right.grid(row=0, column=1, sticky="nsew")
    left.columnconfigure(1, weight=1)
    right.rowconfigure(2, weight=1)

    # Helpers UI
    def log(msg: str):
        q.put(("log", msg))

    def set_busy(busy: bool):
        q.put(("busy", busy))

    # Widgets - Left (inputs)
    pad_y = 8

    ttk.Label(left, text="Arquivo Excel (aba 'Dados')").grid(row=0, column=0, columnspan=2, sticky="w", padx=14, pady=(14, 2))
    excel_entry = ttk.Entry(left, textvariable=excel_var)
    excel_entry.grid(row=1, column=0, sticky="ew", padx=(14, 8), pady=(0, pad_y))
    ttk.Button(left, text="Procurar", style="Neon.TButton",
               command=lambda: excel_var.set(filedialog.askopenfilename(filetypes=[("Excel", "*.xlsx *.xls")]) or excel_var.get())
               ).grid(row=1, column=1, sticky="e", padx=(0, 14), pady=(0, pad_y))

    ttk.Label(left, text="Número do sinistro (Nº Reguladora)").grid(row=2, column=0, columnspan=2, sticky="w", padx=14, pady=(4, 2))
    sin_entry = ttk.Entry(left, textvariable=sinistro_var)
    sin_entry.grid(row=3, column=0, sticky="ew", padx=(14, 8), pady=(0, pad_y))
    ttk.Button(left, text="Validar", style="Neon.TButton").grid(row=3, column=1, sticky="e", padx=(0, 14), pady=(0, pad_y))

    ttk.Label(left, text="Tipo de sinistro (template)").grid(row=4, column=0, columnspan=2, sticky="w", padx=14, pady=(4, 2))
    types_frame = ttk.Frame(left, style="Card.TFrame")
    types_frame.grid(row=5, column=0, columnspan=2, sticky="w", padx=14, pady=(0, pad_y))
    ttk.Radiobutton(types_frame, text="A - Acidente/Avaria/Tombamento/Colisão", variable=tipo_var, value="A").pack(anchor="w")
    ttk.Radiobutton(types_frame, text="R - Roubo/Furto", variable=tipo_var, value="R").pack(anchor="w")
    ttk.Radiobutton(types_frame, text="V - Variação Temperatura", variable=tipo_var, value="V").pack(anchor="w")

    ttk.Label(left, text="Templates (.pptx)").grid(row=6, column=0, columnspan=2, sticky="w", padx=14, pady=(4, 2))
    ttk.Label(left, text="Você pode alterar os caminhos abaixo (um por tipo).", style="Hint.TLabel").grid(
        row=7, column=0, columnspan=2, sticky="w", padx=14, pady=(0, 6)
    )

    def template_row(r, label, var):
        ttk.Label(left, text=label).grid(row=r, column=0, sticky="w", padx=14, pady=(0, 2))
        e = ttk.Entry(left, textvariable=var)
        e.grid(row=r + 1, column=0, sticky="ew", padx=(14, 8), pady=(0, pad_y))
        ttk.Button(left, text="...", style="Neon.TButton",
                   command=lambda: var.set(filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")]) or var.get())
                   ).grid(row=r + 1, column=1, sticky="e", padx=(0, 14), pady=(0, pad_y))

    template_row(8, "Template A", template_a_var)
    template_row(10, "Template R", template_r_var)
    template_row(12, "Template V", template_v_var)

    ttk.Label(left, text="Pasta de saída (PPT + Excel filtrado)").grid(row=14, column=0, columnspan=2, sticky="w", padx=14, pady=(4, 2))
    out_entry = ttk.Entry(left, textvariable=outdir_var)
    out_entry.grid(row=15, column=0, sticky="ew", padx=(14, 8), pady=(0, 14))
    ttk.Button(left, text="Escolher", style="Neon.TButton",
               command=lambda: outdir_var.set(filedialog.askdirectory() or outdir_var.get())
               ).grid(row=15, column=1, sticky="e", padx=(0, 14), pady=(0, 14))

    # Widgets - Right (log + actions)
    ttk.Label(right, text="Execução / Log").grid(row=0, column=0, sticky="w", padx=14, pady=(14, 8))

    progress = ttk.Progressbar(right, mode="indeterminate")
    progress.grid(row=1, column=0, sticky="ew", padx=14, pady=(0, 10))

    log_text = tk.Text(right, height=22, bg="#070a12", fg="#c7d2fe", insertbackground="#c7d2fe",
                       relief="flat", padx=10, pady=10, font=("Consolas", 10))
    log_text.grid(row=2, column=0, sticky="nsew", padx=14, pady=(0, 12))
    log_text.configure(state="disabled")

    def append_log(line: str):
        log_text.configure(state="normal")
        log_text.insert("end", line + "\n")
        log_text.see("end")
        log_text.configure(state="disabled")

    def validate_sinistro():
        try:
            excel_path = excel_var.get().strip()
            if not excel_path or not os.path.exists(excel_path):
                append_log("Erro: arquivo Excel não encontrado.")
                return

            num = sinistro_var.get().strip()
            if not num:
                append_log("Erro: informe o número do sinistro.")
                return

            try:
                num_int = int(num)
            except ValueError:
                append_log("Erro: número do sinistro precisa ser numérico.")
                return

            df = pd.read_excel(excel_path, sheet_name="Dados")
            if "Nº Reguladora" not in df.columns:
                append_log("Erro: coluna 'Nº Reguladora' não existe na aba 'Dados'.")
                return

            ok = (df["Nº Reguladora"] == num_int).any()
            append_log("OK: sinistro encontrado na planilha." if ok else "Não encontrado: sinistro não existe na planilha.")
        except Exception as e:
            append_log(f"Erro ao validar: {e}")

    # vincula botão validar
    for w in left.grid_slaves(row=3, column=1):
        if isinstance(w, ttk.Button):
            w.configure(command=validate_sinistro)

    def get_template_by_tipo(tipo: str) -> str:
        if tipo == "A":
            return template_a_var.get().strip()
        if tipo == "R":
            return template_r_var.get().strip()
        return template_v_var.get().strip()

    def worker_generate():
        try:
            excel_path = excel_var.get().strip()
            if not excel_path or not os.path.exists(excel_path):
                log("Erro: arquivo Excel não encontrado.")
                return

            num_str = sinistro_var.get().strip()
            if not num_str:
                log("Erro: informe o número do sinistro.")
                return
            try:
                num_int = int(num_str)
            except ValueError:
                log("Erro: número do sinistro precisa ser numérico.")
                return

            tipo = tipo_var.get().strip().upper()[:1]
            if tipo not in ("A", "R", "V"):
                log("Erro: tipo de sinistro inválido (A/R/V).")
                return

            template_path = get_template_by_tipo(tipo)
            if not template_path or not os.path.exists(template_path):
                log(f"Erro: template do tipo {tipo} não encontrado: {template_path}")
                return

            outdir = outdir_var.get().strip()
            if not outdir:
                outdir = os.path.dirname(excel_path)

            if not outdir or not os.path.isdir(outdir):
                log("Erro: pasta de saída inválida.")
                return

            set_busy(True)
            out_ppt = criar_apresentacao_sinistro(
                numero_sinistro=num_int,
                excel_source_path=excel_path,
                template_path=template_path,
                output_dir=outdir,
                log_func=log,
            )
            if out_ppt:
                log("Concluído.")
            else:
                log("Falhou: verifique mensagens acima.")
        finally:
            set_busy(False)

    def on_generate_click():
        t = threading.Thread(target=worker_generate, daemon=True)
        t.start()

    actions = ttk.Frame(right, style="Card.TFrame")
    actions.grid(row=3, column=0, sticky="ew", padx=14, pady=(0, 14))
    ttk.Button(actions, text="Gerar PPT", style="Accent.TButton", command=on_generate_click).pack(side="left")
    ttk.Button(actions, text="Limpar Log", style="Neon.TButton", command=lambda: (log_text.configure(state="normal"), log_text.delete("1.0", "end"), log_text.configure(state="disabled"))).pack(side="left", padx=10)

    # Loop de eventos para mensagens de thread
    busy_state = {"on": False}

    def pump_queue():
        try:
            while True:
                kind, payload = q.get_nowait()
                if kind == "log":
                    append_log(payload)
                elif kind == "busy":
                    if payload and not busy_state["on"]:
                        busy_state["on"] = True
                        progress.start(12)
                    if (not payload) and busy_state["on"]:
                        busy_state["on"] = False
                        progress.stop()
        except queue.Empty:
            pass
        root.after(80, pump_queue)

    pump_queue()
    root.mainloop()


# -----------------------------
# CLI (mantém compatibilidade)
# -----------------------------
def run_cli():
    """
    Uso:
      python script.py <numero_sinistro> <tipo(A|R|V)> [excel_path] [template_path] [output_dir]

    Se não passar args -> abre UI.
    """
    if len(sys.argv) < 3:
        print("Args insuficientes. Use a UI (sem args) ou:")
        print("python script.py <numero_sinistro> <tipo(A|R|V)> [excel_path] [template_path] [output_dir]")
        return

    numero_sinistro = int(sys.argv[1])
    tipo = sys.argv[2].strip().upper()[:1]
    excel_path = sys.argv[3] if len(sys.argv) >= 4 else DEFAULT_EXCEL

    if len(sys.argv) >= 5:
        template_path = sys.argv[4]
    else:
        template_path = DEFAULT_TEMPLATE_A if tipo == "A" else DEFAULT_TEMPLATE_R if tipo == "R" else DEFAULT_TEMPLATE_V

    output_dir = sys.argv[5] if len(sys.argv) >= 6 else os.path.dirname(excel_path)

    def log(msg):
        print(msg)

    criar_apresentacao_sinistro(
        numero_sinistro=numero_sinistro,
        excel_source_path=excel_path,
        template_path=template_path,
        output_dir=output_dir,
        log_func=log,
    )


if __name__ == "__main__":
    if len(sys.argv) > 1:
        run_cli()
    else:
        run_ui()