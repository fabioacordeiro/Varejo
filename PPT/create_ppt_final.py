# pip install python-pptx
# pip install pandas
# pip install openpyxl


import pandas as pd
from pptx import Presentation
import os
import re
import sys
from datetime import datetime

def criar_apresentacao_sinistro(numero_sinistro):
    # Carregar os dados do Excel
    excel_file = 'C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\SINISTROS1.xlsx'
    df = pd.read_excel(excel_file)
    
    # Filtrar para o sinistro especificado (na coluna Status_Cordeiro)
    sinistro_df = df[df['Status_Cordeiro'] == numero_sinistro]
    
    if sinistro_df.empty:
        print(f"Sinistro {numero_sinistro} não encontrado.")
        return False
    
    # Extrair os dados do sinistro
    sinistro_data = sinistro_df.iloc[0]
    
    # Formatar o valor do prejuízo
    prejuizo = sinistro_data['Unnamed: 10']
    if pd.notna(prejuizo) and isinstance(prejuizo, (int, float)):
        prejuizo_formatado = f"R$ {prejuizo:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        prejuizo_formatado = str(prejuizo)
    
    # Formatar a data (assumindo formato YYYY_M)
    data_str = str(sinistro_data['Unnamed: 3'])
    if re.match(r'\d{4}_\d{1,2}', data_str):
        ano, mes = data_str.split('_')
        mes_int = int(mes)
        nomes_meses = ["Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho", 
                      "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"]
        if 1 <= mes_int <= 12:
            data_formatada = f"{nomes_meses[mes_int-1]} de {ano}"
        else:
            data_formatada = data_str
    else:
        data_formatada = data_str
    
    # Carregar o template PPT
    template_path = 'C:\\Fabio\\Desenvolvimento\\Varejo\\ppt\\Sinistro_65329.pptx'
    output_path = f'C:\\Fabio\\Desenvolvimento\\Varejo\\ppt\\Sinistro_{numero_sinistro}_Gerado.pptx'
    
    # Criar uma cópia do template
    prs = Presentation(template_path)
    
    # Mapear os dados do sinistro para os campos do template
    sinistro_info = {
        'Número do Sinistro': str(numero_sinistro),
        'Status': str(sinistro_data['EM ANDAMENTO']),
        'Data': data_formatada,
        'Causa': str(sinistro_data['Unnamed: 4']),
        'Transportador': str(sinistro_data['Unnamed: 5']),
        'UF Origem': str(sinistro_data['Unnamed: 6']),
        'Cidade Origem': str(sinistro_data['Unnamed: 7']),
        'UF Destino': str(sinistro_data['Unnamed: 8']),
        'Cidade Destino': str(sinistro_data['Unnamed: 9']),
        'Prejuízo Apurado': prejuizo_formatado,
        'Descrição': str(sinistro_data['Unnamed: 11'])
    }
    
    # Função para substituir texto mantendo a formatação
    def replace_text_in_shape(shape, replacements):
        if not hasattr(shape, "text_frame"):
            return
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, value)
    
    # Atualizar os slides com os dados do sinistro
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                replace_text_in_shape(shape, sinistro_info)
            
            # Processar tabelas, se houver
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_shape(cell, sinistro_info)
    
    # Salvar a apresentação
    prs.save(output_path)
    print(f"Apresentação gerada com sucesso: {output_path}")
    return True

if __name__ == "__main__":
    # Se um número de sinistro for fornecido como argumento, use-o
    if len(sys.argv) > 1:
        try:
            numero_sinistro = int(sys.argv[1])
            criar_apresentacao_sinistro(numero_sinistro)
        except ValueError:
            print("Erro: O número do sinistro deve ser um valor numérico.")
    else:
        # Caso contrário, use o número de sinistro padrão (65329)
        criar_apresentacao_sinistro(60981)

