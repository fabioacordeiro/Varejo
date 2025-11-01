# pip install python-pptx
# pip install pandas
# pip install openpyxl
# Este script depende dos arquivos de Excel para filtro e PPT (Template) para criar um novo com os 
# novos dados.
# Criado por Fábio A Cordeiro em 24/08/2025 - 
# Created by: Fábio A Cordeiro (fabioacordeiro@yahoo.com.br - 11-95456-7877)

import pandas as pd
from pptx import Presentation
import os
import re
import sys
from datetime import datetime, timedelta
import numpy as np

# Caminho do arquivo
file_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\SINISTROS1.xlsx"

# Ler a planilha somente a pasta Dados
df = pd.read_excel(file_path, sheet_name="Dados")

# Exibir na tela os nomes das colunas encontradas
print("Colunas encontradas:")
print(df.columns.tolist())

print("Exibir Dataframe")
print(df)

# >>> Tipos das colunas (incluído a seu pedido)
print("dtypes do DataFrame:")
print(df.dtypes)

# Definir a coluna de filtro
coluna_filtro = "Nº Reguladora"

# Exemplo: filtrar somente registros onde há número de reguladora válido
df_filtrado = df[df[coluna_filtro].notna()]
# Filtrar para o sinistro especificado (na coluna Nº Reguladora)

num_sinistro = int(input('Digite o número do sinistro para criação do ppt :'))
sinistro_df = df[df['Nº Reguladora'] == num_sinistro]

# Exibir algumas linhas
print(df)
print("\nDados filtrados:")
print(df_filtrado.head())

# Se quiser salvar em novo Excel com os dados filtrados
sinistro_df.to_excel("C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\SINISTROS_FILTRADO.xlsx", index=False)

# -------------------------------------------------------------
# Funções utilitárias de data/hora (NOVA)
# -------------------------------------------------------------
def format_datetime(value):
    """
    Normaliza valores de data/hora e devolve string no formato "%Y-%m-%d %H:%M".
    Aceita: datetime, pandas.Timestamp, string, número serial do Excel, NaN/None.
    """
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return ""
    # datetime / Timestamp
    if isinstance(value, (datetime, pd.Timestamp)):
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y %H:%M")
    # número (serial do Excel)
    if isinstance(value, (int, float, np.floating)):
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y %H:%M")
        except Exception:
            return ""
    # string
    if isinstance(value, str):
        s = value.strip()
        for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y %H:%M")
            except ValueError:
                continue
        # fallback: mantém original (útil para depurar o que veio)
        return value
    # outros tipos: devolve string
    return str(value)

def criar_apresentacao_sinistro(numero_sinistro=None, excel_file=None, template_path=None):
    """
    Cria uma apresentação PowerPoint baseada em um template e dados de uma planilha Excel.
    
    Args:
        numero_sinistro: Número do sinistro para filtrar 
        excel_file: Caminho para o arquivo Excel 
        template_path: Caminho para o template PPT 
    
    Returns:
        str: Caminho do arquivo PPT gerado ou None em caso de erro
    """
    # Caminho da planilha
    if excel_file is None:
        excel_file = 'C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\SINISTROS_FILTRADO.xlsx'
    
    if template_path is None:
        template_path = 'C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\Template.pptx'
    
    # Ler o arquivo Excel e Carregar os dados
    try:
        df = pd.read_excel(excel_file)
        print(f"Dados carregados do Excel: {len(df)} linhas")
    except Exception as e:
        print(f"Erro ao ler o arquivo Excel: {e}")
        return None
    
    # Verificar se há dados
    if df.empty:
        print("A planilha Excel está vazia.")
        return None
    
    # Se um número de sinistro for fornecido, filtrar os dados
    if numero_sinistro is not None:
        if "Nº Reguladora" in df.columns:
            df = df[df["Nº Reguladora"] == numero_sinistro]
            if df.empty:
                print(f"Sinistro {numero_sinistro} não encontrado na planilha.")
                return None
    
    # Extrair os dados do sinistro (primeira linha da planilha)
    sinistro_data = df.iloc[0]

    # (Opcional) Ver dtypes do DF carregado aqui também:
    print("-"*58)
    print("dtypes (após carregar excel_file):")
    print(df.dtypes)
    print("-"*58)

    # Obter o número do sinistro dos dados
    if numero_sinistro is None:
        if "Nº Reguladora" in sinistro_data:
            numero_sinistro = sinistro_data["Nº Reguladora"]
        else:
            numero_sinistro = "Desconhecido"
    
    print(f"Processando sinistro número: {numero_sinistro}")
    
    # Definir o nome e o caminho de gravação dos dados
    output_path = f'C:\\Fabio\\Desenvolvimento\\Varejo\\PPT\\Sinistro_{numero_sinistro}_Final_v2.pptx'
    
    # Formatar o valor do prejuízo
    prejuizo = sinistro_data.get("Prejuizo Apurado", "N/A")
    if pd.notna(prejuizo) and isinstance(prejuizo, (int, float, np.floating)):
        prejuizo_formatado = f"R$ {prejuizo:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        prejuizo_formatado = str(prejuizo)
    
    # Formatar o valor do embarque
    valor_embarque = sinistro_data.get("Valor do Embarque", "N/A")
    if pd.notna(valor_embarque) and isinstance(valor_embarque, (int, float, np.floating)):
        valor_embarque_formatado = f"R$ {valor_embarque:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        valor_embarque_formatado = str(valor_embarque)
    
    # Formatar a data do sinistro (mantido dd/mm/aaaa, como no seu original)
    data_sinistro = sinistro_data.get("Data do Sinistro", "N/A")
    if isinstance(data_sinistro, (datetime, pd.Timestamp)):
        data_formatada = (data_sinistro.to_pydatetime() if isinstance(data_sinistro, pd.Timestamp) else data_sinistro).strftime("%d/%m/%Y")
    elif isinstance(data_sinistro, str):
        # tentativa simples
        try:
            data_formatada = datetime.strptime(data_sinistro.strip(), "%d/%m/%Y").strftime("%d/%m/%Y")
        except Exception:
            data_formatada = str(data_sinistro)
    else:
        data_formatada = str(data_sinistro)

    # Carregar o template PPT
    try:
        prs = Presentation(template_path)
        print(f"Template PPT carregado: {template_path}")
    except Exception as e:
        print(f"Erro ao carregar o template PPT: {e}")
        return None
    
    # Extrair informações do template para substituição
    template_info = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and "SINISTRO" in shape.text_frame.text:
                match = re.search(r'SINISTRO (\d+\.?\d*)', shape.text_frame.text)
                if match:
                    template_info["numero_sinistro"] = match.group(1)
                
                match = re.search(r'– ([^–]+) –', shape.text_frame.text)
                if match:
                    template_info["cidade_origem"] = match.group(1).strip()
                
                match = re.search(r'\(([^)]+)\)', shape.text_frame.text)
                if match:
                    template_info["destino_template"] = match.group(1).strip()
    
    print(f"Informações extraídas do template: {template_info}")
    
    # Criar título para o slide
    causa = str(sinistro_data.get("Causa Final", "")).strip()
    cidade_origem = str(sinistro_data.get("Origem Ajustado", "")).strip()
    cidade_destino = str(sinistro_data.get("Cidade - Destino", "")).strip()
    uf_origem = str(sinistro_data.get("UF - Origem", "")).strip()
    uf_destino = str(sinistro_data.get("UF - Destino", "")).strip()
    
    titulo_slide = f"SINISTRO {numero_sinistro} – {causa} – {cidade_origem} – ({cidade_destino})"
    
    # Mapear os dados do sinistro para os campos do template
    sinistro_info = {
        # Substituições específicas do template
        template_info.get("numero_sinistro", "65.329"): str(numero_sinistro),
        template_info.get("cidade_origem", "ITAPECERICA DA SERRA"): cidade_origem,
        template_info.get("Complemento_info", " - "): cidade_destino,
        
        # Substituições gerais
        "SINISTRO ": f"SINISTRO {numero_sinistro}",
        #"Número": str(numero_sinistro),
        "CAUSA": causa,
        "Info_CidadeOrigem": cidade_origem,
        "Info_Cid_Origem": str(sinistro_data.get("Cidade Origem", "")),
        "Info_Uf_Origem": uf_origem,
        "Info_Destino": cidade_destino,
        "Info_Cid_Destino": str(sinistro_data.get("Complemento_info", "")),
        "Info_Transp": str(sinistro_data.get("Transportador", "")),
        "Info_Placa": str(sinistro_data.get("Placa", "")),
        "Info_mot": str(sinistro_data.get("Motorista", "")),
        "Info_VlCarga": f"R$ {valor_embarque_formatado}",
        "Info_Prejuizo": f"R$ {prejuizo_formatado}",
        "Info_DT_Sinistro": data_formatada,
        "Info_Carga": str(sinistro_data.get("N_Carga", "")),
        "Info_Qte_Cargas": str(sinistro_data.get("QTDE_CARGAS_MOT", "")),
        "Info_Desc": str(sinistro_data.get("Ação", "")),

        # >>> DATAS FORMATADAS (YYYY-mm-dd HH:MM) <<<
        "Info_Doca":           format_datetime(sinistro_data.get("ENCOSTA_EM_DOCA", "")),
        "Info_Inicio_Carreg":  format_datetime(sinistro_data.get("INICIO_CARREGAMENTO", "")),
        "Info_Fim_Carreg":     format_datetime(sinistro_data.get("FIM_CARREGAMENTO", "")),
        "Info_Emissao_NF":     format_datetime(sinistro_data.get("EMISSAO_NF", "")),
        "Info_Inicio_Viagem":  format_datetime(sinistro_data.get("INICIO_VIAGEM", "")),
        "Info_Chegada":        format_datetime(sinistro_data.get("CHEGADA_EM_LOJA", "")),
    }
    
    # Adicionar o texto completo da descrição para substituição
    descricao = str(sinistro_data.get("Ação", ""))
    if descricao:
        # Encontrar o texto da descrição no template para substituir
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text and "Veículo saiu carregado" in cell.text_frame.text:
                                sinistro_info[cell.text_frame.text] = descricao
                            elif cell.text_frame.text and "Resumo Ocorrência" in cell.text_frame.text:
                                # Encontrar a célula adjacente para substituir o texto
                                if len(row.cells) > 1:
                                    for adj_cell in row.cells:
                                        if adj_cell != cell and adj_cell.text_frame.text:
                                            sinistro_info[adj_cell.text_frame.text] = descricao
    
    print(f"Mapeamento de substituição criado com {len(sinistro_info)} itens")
    
    # Função para substituir texto mantendo a formatação
    def replace_text_in_shape(shape, replacements):
        if not hasattr(shape, "text_frame"):
            return False
        
        replaced = False
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                new_text = original_text
                
                # Substituir textos exatos
                for key, value in replacements.items():
                    if key in new_text:
                        new_text = new_text.replace(key, str(value))
                        replaced = True
                
                # Substituir o título completo do slide
                if "SINISTRO" in new_text and "BACKOFFICE" in new_text:
                    new_text = f"{titulo_slide}\nBACKOFFICE - SINISTROS"
                    replaced = True
                
                if new_text != original_text:
                    run.text = new_text
        
        return replaced
    
    # Atualizar os slides com os dados do sinistro
    total_replacements = 0
    for slide_idx, slide in enumerate(prs.slides):
        slide_replacements = 0
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if replace_text_in_shape(shape, sinistro_info):
                    slide_replacements += 1
            
            # Processar tabelas, se houver
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_text_in_shape(cell, sinistro_info):
                            slide_replacements += 1
        
        print(f"Slide {slide_idx+1}: {slide_replacements} substituições realizadas")
        total_replacements += slide_replacements
    
    print(f"Total de substituições realizadas: {total_replacements}")
    
    # Salvar a apresentação
    try:
        prs.save(output_path)
        print(f"Apresentação gerada com sucesso: {output_path}")
        return output_path
    except Exception as e:
        print(f"Erro ao salvar a apresentação: {e}")
        return None

if __name__ == "__main__":
    # Se um número de sinistro for fornecido como argumento, use-o
    if len(sys.argv) > 1:
        try:
            numero_sinistro = int(sys.argv[1])
            criar_apresentacao_sinistro(numero_sinistro)
        except ValueError:
            print("Erro: O número do sinistro deve ser um valor numérico.")
    else:
        # Caso contrário, use os dados da primeira linha da planilha
        criar_apresentacao_sinistro()