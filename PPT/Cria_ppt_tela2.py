# pip install python-pptx pandas openpyxl numpy

import os
import re
import sys
from datetime import datetime, timedelta
import numpy as np
import pandas as pd
from pptx import Presentation

import tkinter as tk
from tkinter import ttk, filedialog, messagebox


# =============================================================================
# 1) FUNÇÕES "SAFE" (CORREÇÃO DEFINITIVA DO NaTType does not support strftime)
# =============================================================================

def safe_isna(v) -> bool:
    """True para None, NaN, NaT, etc."""
    try:
        return v is None or pd.isna(v)
    except Exception:
        return v is None


def safe_date(value) -> str:
    """Retorna dd/mm/YYYY. Nunca quebra com NaT."""
    if safe_isna(value):
        return ""

    if isinstance(value, (datetime, pd.Timestamp)):
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y")

    if isinstance(value, (int, float, np.floating)):
        # serial Excel
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y")
        except Exception:
            return ""

    if isinstance(value, str):
        s = value.strip()
        if not s or s.lower() in ("nan", "nat", "n/a", "none", "-"):
            return ""
        for fmt in ("%d/%m/%Y", "%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y")
            except ValueError:
                pass
        return s  # fallback

    return str(value)


def safe_dt(value) -> str:
    """Retorna dd/mm/YYYY HH:MM. Nunca quebra com NaT."""
    if safe_isna(value):
        return ""

    if isinstance(value, (datetime, pd.Timestamp)):
        dt = value.to_pydatetime() if isinstance(value, pd.Timestamp) else value
        return dt.strftime("%d/%m/%Y %H:%M")

    if isinstance(value, (int, float, np.floating)):
        # serial Excel
        try:
            dt = datetime(1899, 12, 30) + timedelta(days=float(value))
            return dt.strftime("%d/%m/%Y %H:%M")
        except Exception:
            return ""

    if isinstance(value, str):
        s = value.strip()
        if not s or s.lower() in ("nan", "nat", "n/a", "none", "-"):
            return ""
        for fmt in ("%d/%m/%Y %H:%M:%S", "%d/%m/%Y %H:%M", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
            try:
                return datetime.strptime(s, fmt).strftime("%d/%m/%Y %H:%M")
            except ValueError:
                pass
        return s

    return str(value)


def format_brl_currency(value) -> str:
    """Formata número para BRL sem quebrar com NaN."""
    if safe_isna(value):
        return ""
    if isinstance(value, (int, float, np.floating)):
        s = f"{float(value):,.2f}"
        return s.replace(",", "X").replace(".", ",").replace("X", ".")
    return str(value)


# =============================================================================
# 2) CORE: GERAÇÃO DO PPT (ROBUSTO)
# =============================================================================

def extract_template_info(prs: Presentation) -> dict:
    """
    Tenta extrair chaves do template a partir de um shape que contém 'SINISTRO'.
    Mantive sua lógica, mas sem depender disso para funcionar.
    """
    template_info = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                txt = shape.text_frame.text
                if "SINISTRO" in txt:
                    m = re.search(r"SINISTRO\s+(\d+\.?\d*)", txt)
                    if m:
                        template_info["numero_sinistro"] = m.group(1)

                    m = re.search(r"–\s*([^–]+)\s*–", txt)
                    if m:
                        template_info["cidade_origem"] = m.group(1).strip()

                    m = re.search(r"\(([^)]+)\)", txt)
                    if m:
                        template_info["destino_template"] = m.group(1).strip()
    return template_info


def replace_text_in_shape(shape, replacements, titulo_slide: str) -> bool:
    """
    Substitui texto no shape mantendo formatação (run por run).
    Retorna True se substituiu algo.
    """
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return False

    replaced_any = False

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text
            new_text = original

            # substituições gerais
            for k, v in replacements.items():
                if k and k in new_text:
                    new_text = new_text.replace(k, str(v))
                    replaced_any = True

            # título do slide (sua regra)
            if "SINISTRO" in new_text and "BACKOFFICE" in new_text:
                new_text = f"{titulo_slide}\nBACKOFFICE - SINISTROS"
                replaced_any = True

            if new_text != original:
                run.text = new_text

    return replaced_any


def criar_apresentacao_sinistro(
    numero_sinistro: int,
    excel_file: str,
    template_path: str,
    output_dir: str,
    sheet_name: str = "Dados",
    debug: bool = False
) -> str:
    """
    Lê Excel, filtra sinistro, carrega template PPT e faz substituições.
    Retorna caminho do PPT gerado.
    """

    # 1) Ler Excel
    df = pd.read_excel(excel_file, sheet_name=sheet_name)

    if df.empty:
        raise ValueError("Planilha está vazia.")

    if "Nº Reguladora" not in df.columns:
        raise ValueError("Coluna obrigatória não encontrada: 'Nº Reguladora'.")

    # 2) Filtrar sinistro
    sinistro_df = df[df["Nº Reguladora"] == numero_sinistro]
    if sinistro_df.empty:
        raise ValueError(f"Sinistro {numero_sinistro} não encontrado.")

    sinistro_data = sinistro_df.iloc[0]

    # Debug: ver campos de data que podem virar NaT
    if debug:
        campos_dt = [
            "Data do Sinistro",
            "ENCOSTA_EM_DOCA",
            "INICIO_CARREGAMENTO",
            "FIM_CARREGAMENTO",
            "EMISSAO_NF",
            "INICIO_VIAGEM",
            "CHEGADA_EM_LOJA",
        ]
        print("\n" + "=" * 80)
        print(f"DEBUG (sinistro {numero_sinistro}) - campos datetime")
        for c in campos_dt:
            if c in sinistro_df.columns:
                v = sinistro_data.get(c)
                print(f"- {c}: {repr(v)} | type={type(v).__name__} | isna={safe_isna(v)}")
            else:
                print(f"- {c}: (coluna não existe)")
        print("=" * 80 + "\n")

    # 3) Carregar template
    prs = Presentation(template_path)
    template_info = extract_template_info(prs)

    # 4) Montar dados
    causa = str(sinistro_data.get("Causa Final", "")).strip()
    cidade_origem = str(sinistro_data.get("Origem Ajustado", "")).strip()
    cidade_destino = str(sinistro_data.get("Cidade - Destino", "")).strip()
    uf_origem = str(sinistro_data.get("UF - Origem", "")).strip()
    uf_destino = str(sinistro_data.get("UF - Destino", "")).strip()

    titulo_slide = f"SINISTRO {numero_sinistro} – {causa} – {cidade_origem} – ({cidade_destino})"

    prejuizo = sinistro_data.get("Prejuizo Apurado", None)
    valor_embarque = sinistro_data.get("Valor do Embarque", None)

    prejuizo_formatado = format_brl_currency(prejuizo)
    valor_embarque_formatado = format_brl_currency(valor_embarque)

    data_formatada = safe_date(sinistro_data.get("Data do Sinistro", None))

    # OBS: no seu código original você faz f"R$ {valor_embarque_formatado}" e ainda tinha "R$ " duplicado.
    # Aqui eu deixo padronizado: já retorna "1.234,56" e eu adiciono "R$ " 1 vez.
    sinistro_info = {
        # Substituições específicas do template (se existir)
        template_info.get("numero_sinistro", "65.329"): str(numero_sinistro),
        template_info.get("cidade_origem", "ITAPECERICA DA SERRA"): cidade_origem,
        template_info.get("Complemento_info", " - "): cidade_destino,

        # Chaves gerais do seu template
        "SINISTRO ": f"SINISTRO {numero_sinistro}",
        "CAUSA": causa,
        "Info_CidadeOrigem": cidade_origem,
        "Info_Cid_Origem": str(sinistro_data.get("Cidade Origem", "")),
        "Info_Uf_Origem": uf_origem,
        "Info_Destino": cidade_destino,
        "Info_Cid_Destino": str(sinistro_data.get("Complemento_info", "")),
        "Info_Transp": str(sinistro_data.get("Transportador", "")),
        "Info_Placa": str(sinistro_data.get("Placa", "")),
        "Info_mot": str(sinistro_data.get("Motorista", "")),
        "Info_VlCarga": f"R$ {valor_embarque_formatado}" if valor_embarque_formatado else "",
        "Info_Prejuizo": f"R$ {prejuizo_formatado}" if prejuizo_formatado else "",
        "Info_DT_Sinistro": data_formatada,
        "Info_Carga": str(sinistro_data.get("N_Carga", "")),
        "Info_Qte_Cargas": str(sinistro_data.get("QTDE_CARGAS_MOT", "")),
        "Info_Desc": str(sinistro_data.get("Ação", "")),

        # DATETIME seguros (não quebram com NaT)
        "Info_Doca":           safe_dt(sinistro_data.get("ENCOSTA_EM_DOCA", None)),
        "Info_Inicio_Carreg":  safe_dt(sinistro_data.get("INICIO_CARREGAMENTO", None)),
        "Info_Fim_Carreg":     safe_dt(sinistro_data.get("FIM_CARREGAMENTO", None)),
        "Info_Emissao_NF":     safe_dt(sinistro_data.get("EMISSAO_NF", None)),
        "Info_Inicio_Viagem":  safe_dt(sinistro_data.get("INICIO_VIAGEM", None)),
        "Info_Chegada":        safe_dt(sinistro_data.get("CHEGADA_EM_LOJA", None)),
    }

    # 5) Se quiser substituir descrição longa em tabelas (sua lógica adaptada)
    descricao = str(sinistro_data.get("Ação", "") or "")
    if descricao:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            txt = cell.text_frame.text if cell.text_frame else ""
                            if txt and "Veículo saiu carregado" in txt:
                                sinistro_info[txt] = descricao

    # 6) Aplicar substituições
    total_replacements = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if replace_text_in_shape(shape, sinistro_info, titulo_slide):
                    total_replacements += 1

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_text_in_shape(cell, sinistro_info, titulo_slide):
                            total_replacements += 1

    # 7) Salvar
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"Sinistro_{numero_sinistro}_Final.pptx")
    prs.save(output_path)

    if debug:
        print(f"Total de shapes/cells com substituição: {total_replacements}")
        print(f"PPT gerado: {output_path}")

    return output_path


# =============================================================================
# 3) UI Tkinter (3 abas, visual escuro/futurista)
# =============================================================================

class FuturisticApp(tk.Tk):
    def __init__(self):
        super().__init__()

        self.title("Gerador de PPT - Sinistros (Backoffice)")
        self.geometry("980x620")
        self.minsize(920, 560)

        # Tema escuro
        self.bg = "#0b0f1a"         # fundo principal
        self.panel = "#101a2b"      # painéis
        self.panel2 = "#0f172a"
        self.text = "#e5e7eb"       # texto
        self.muted = "#9ca3af"      # texto secundário
        self.accent = "#6E3AB8"     # roxo (seu padrão)
        self.ok = "#22c55e"
        self.warn = "#f59e0b"
        self.err = "#ef4444"

        self.configure(bg=self.bg)

        # Estilo ttk
        style = ttk.Style(self)
        style.theme_use("clam")

        style.configure("TNotebook", background=self.bg, borderwidth=0)
        style.configure("TNotebook.Tab", background=self.panel2, foreground=self.text, padding=(14, 8))
        style.map("TNotebook.Tab",
                  background=[("selected", self.panel)],
                  foreground=[("selected", "white")])

        style.configure("TFrame", background=self.bg)
        style.configure("Card.TFrame", background=self.panel, relief="flat")

        style.configure("TLabel", background=self.panel, foreground=self.text)
        style.configure("Muted.TLabel", background=self.panel, foreground=self.muted)

        style.configure("TEntry", fieldbackground="#0b1220", foreground=self.text, insertcolor=self.text)
        style.configure("TCombobox", fieldbackground="#0b1220", foreground=self.text)

        style.configure("Accent.TButton", background=self.accent, foreground="white", padding=(12, 8))
        style.map("Accent.TButton",
                  background=[("active", "#7c4ad1")])

        style.configure("Ghost.TButton", background=self.panel2, foreground=self.text, padding=(12, 8))
        style.map("Ghost.TButton",
                  background=[("active", "#1b2740")])

        # Variáveis
        self.excel_path = tk.StringVar(value=r"C:\Fabio\Desenvolvimento\Varejo\Sinistro\SINISTROS1.xlsx")
        self.sheet_name = tk.StringVar(value="Dados")
        self.output_dir = tk.StringVar(value=r"C:\Fabio\Desenvolvimento\Varejo\PPT")

        self.tipo_sinistro = tk.StringVar(value="A")  # A/R/V
        self.template_a = tk.StringVar(value=r"C:\Fabio\Desenvolvimento\Varejo\PPT\TemplateA.pptx")
        self.template_r = tk.StringVar(value=r"C:\Fabio\Desenvolvimento\Varejo\PPT\TemplateR.pptx")
        self.template_v = tk.StringVar(value=r"C:\Fabio\Desenvolvimento\Varejo\PPT\TemplateV.pptx")

        self.numero_sinistro = tk.StringVar(value="")
        self.debug_mode = tk.BooleanVar(value=False)

        self._build_ui()

    def _build_ui(self):
        nb = ttk.Notebook(self)
        nb.pack(fill="both", expand=True, padx=14, pady=14)

        # Aba 1: Config
        tab1 = ttk.Frame(nb, style="TFrame")
        nb.add(tab1, text="Configurações")

        card1 = ttk.Frame(tab1, style="Card.TFrame")
        card1.pack(fill="x", padx=10, pady=10)

        self._row_filepicker(card1, "Excel (SINISTROS1.xlsx)", self.excel_path, self._pick_excel)
        self._row_entry(card1, "Aba (sheet_name)", self.sheet_name)
        self._row_filepicker(card1, "Pasta de saída (PPT)", self.output_dir, self._pick_output_dir)

        # Aba 2: Templates
        tab2 = ttk.Frame(nb, style="TFrame")
        nb.add(tab2, text="Templates")

        card2 = ttk.Frame(tab2, style="Card.TFrame")
        card2.pack(fill="x", padx=10, pady=10)

        self._row_filepicker(card2, "Template A (Acidente/Avaria/...)", self.template_a, self._pick_template_a)
        self._row_filepicker(card2, "Template R (Roubo/Furto)", self.template_r, self._pick_template_r)
        self._row_filepicker(card2, "Template V (Variação Temperatura)", self.template_v, self._pick_template_v)

        # Aba 3: Gerar
        tab3 = ttk.Frame(nb, style="TFrame")
        nb.add(tab3, text="Gerar PPT")

        card3 = ttk.Frame(tab3, style="Card.TFrame")
        card3.pack(fill="x", padx=10, pady=10)

        # tipo sinistro
        row = ttk.Frame(card3, style="Card.TFrame")
        row.pack(fill="x", padx=14, pady=(14, 6))
        ttk.Label(row, text="Tipo de Sinistro", style="TLabel").pack(side="left")
        cb = ttk.Combobox(row, textvariable=self.tipo_sinistro, values=["A", "R", "V"], width=10, state="readonly")
        cb.pack(side="left", padx=12)

        # numero sinistro
        self._row_entry(card3, "Nº do Sinistro (Nº Reguladora)", self.numero_sinistro)

        # debug
        rowd = ttk.Frame(card3, style="Card.TFrame")
        rowd.pack(fill="x", padx=14, pady=(6, 6))
        chk = tk.Checkbutton(
            rowd, text="Debug (imprimir campos de data no console)",
            variable=self.debug_mode,
            bg=self.panel, fg=self.text,
            selectcolor=self.panel2, activebackground=self.panel, activeforeground="white"
        )
        chk.pack(side="left")

        # botões
        rowb = ttk.Frame(card3, style="Card.TFrame")
        rowb.pack(fill="x", padx=14, pady=(12, 14))

        btn_generate = ttk.Button(rowb, text="Gerar PPT", style="Accent.TButton", command=self._generate)
        btn_generate.pack(side="left")

        btn_clear = ttk.Button(rowb, text="Limpar", style="Ghost.TButton", command=self._clear)
        btn_clear.pack(side="left", padx=10)

        # log area
        cardlog = ttk.Frame(tab3, style="Card.TFrame")
        cardlog.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        self.log = tk.Text(
            cardlog, height=12, wrap="word",
            bg="#050814", fg=self.text, insertbackground=self.text,
            relief="flat", padx=12, pady=10
        )
        self.log.pack(fill="both", expand=True, padx=10, pady=10)
        self._log("Pronto. Informe o sinistro e clique em Gerar PPT.")

    def _row_entry(self, parent, label, var):
        row = ttk.Frame(parent, style="Card.TFrame")
        row.pack(fill="x", padx=14, pady=6)

        ttk.Label(row, text=label, style="TLabel").pack(side="left")
        ent = ttk.Entry(row, textvariable=var, width=70)
        ent.pack(side="left", padx=12, fill="x", expand=True)

    def _row_filepicker(self, parent, label, var, cmd_pick):
        row = ttk.Frame(parent, style="Card.TFrame")
        row.pack(fill="x", padx=14, pady=6)

        ttk.Label(row, text=label, style="TLabel").pack(side="left")
        ent = ttk.Entry(row, textvariable=var, width=70)
        ent.pack(side="left", padx=12, fill="x", expand=True)
        ttk.Button(row, text="Procurar...", style="Ghost.TButton", command=cmd_pick).pack(side="left")

    def _pick_excel(self):
        p = filedialog.askopenfilename(title="Selecione o Excel", filetypes=[("Excel", "*.xlsx *.xls")])
        if p:
            self.excel_path.set(p)

    def _pick_output_dir(self):
        p = filedialog.askdirectory(title="Selecione a pasta de saída")
        if p:
            self.output_dir.set(p)

    def _pick_template_a(self):
        p = filedialog.askopenfilename(title="Selecione o Template A", filetypes=[("PowerPoint", "*.pptx")])
        if p:
            self.template_a.set(p)

    def _pick_template_r(self):
        p = filedialog.askopenfilename(title="Selecione o Template R", filetypes=[("PowerPoint", "*.pptx")])
        if p:
            self.template_r.set(p)

    def _pick_template_v(self):
        p = filedialog.askopenfilename(title="Selecione o Template V", filetypes=[("PowerPoint", "*.pptx")])
        if p:
            self.template_v.set(p)

    def _clear(self):
        self.numero_sinistro.set("")
        self._log("Campos limpos.")

    def _get_template_by_type(self) -> str:
        t = self.tipo_sinistro.get().strip().upper()
        if t == "A":
            return self.template_a.get()
        if t == "R":
            return self.template_r.get()
        if t == "V":
            return self.template_v.get()
        return self.template_a.get()

    def _log(self, msg: str):
        ts = datetime.now().strftime("%H:%M:%S")
        self.log.insert("end", f"[{ts}] {msg}\n")
        self.log.see("end")

    def _validate_paths(self):
        excel = self.excel_path.get().strip()
        if not os.path.isfile(excel):
            raise ValueError("Arquivo Excel não encontrado.")

        template = self._get_template_by_type()
        if not os.path.isfile(template):
            raise ValueError("Template PPTX não encontrado para o tipo selecionado.")

        outdir = self.output_dir.get().strip()
        if not outdir:
            raise ValueError("Pasta de saída inválida.")
        return excel, template, outdir

    def _generate(self):
        try:
            excel, template, outdir = self._validate_paths()

            ns = self.numero_sinistro.get().strip()
            if not ns.isdigit():
                raise ValueError("Informe um Nº de sinistro numérico (Nº Reguladora).")
            numero = int(ns)

            self._log(f"Iniciando geração: sinistro={numero} | tipo={self.tipo_sinistro.get()}")

            ppt_path = criar_apresentacao_sinistro(
                numero_sinistro=numero,
                excel_file=excel,
                template_path=template,
                output_dir=outdir,
                sheet_name=self.sheet_name.get().strip() or "Dados",
                debug=self.debug_mode.get()
            )

            self._log(f"OK: PPT gerado com sucesso.")
            self._log(f"Arquivo: {ppt_path}")
            messagebox.showinfo("Sucesso", "PPT gerado com sucesso!")

        except Exception as e:
            self._log(f"ERRO: {e}")
            messagebox.showerror("Erro", str(e))


if __name__ == "__main__":
    app = FuturisticApp()
    app.mainloop()