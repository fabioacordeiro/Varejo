import pandas as pd
from pptx import Presentation
import os
import re
import sys
from datetime import datetime

def criar_apresentacao_sinistro(numero_sinistro):
    # Caminho do arquivo Excel
    excel_file = "SINISTROS1.xlsx"
    
    # Carregar os dados do Excel, usando a segunda linha como cabeçalho
    #df = pd.read_excel(excel_file, sheet_name="Dados", header=1)
    df = pd.read_excel(excel_file, sheet_name="Dados")
    
    # Filtrar para o sinistro especificado na coluna 'Nº Reguladora'
    sinistro_df = df[df["Nº Reguladora"] == numero_sinistro]
    
    if sinistro_df.empty:
        print(f"Sinistro {numero_sinistro} não encontrado na coluna 'Nº Reguladora'.")
        return False
    
    # Extrair os dados do sinistro
    sinistro_data = sinistro_df.iloc[0]
    
    # Formatar o valor do prejuízo
    prejuizo = sinistro_data["Prejuizo Apurado"]
    if pd.notna(prejuizo) and isinstance(prejuizo, (int, float)):
        prejuizo_formatado = f"R$ {prejuizo:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        prejuizo_formatado = str(prejuizo)
    
    # Formatar a data (assumindo formato YYYY_M ou Data do Sinistro)
    data_str = str(sinistro_data["Data do Sinistro"])
    data_formatada = data_str # Usar a string original por padrão
    
    # Tentar formatar se for um objeto datetime ou string YYYY-MM-DD
    if isinstance(sinistro_data["Data do Sinistro"], datetime):
        data_formatada = sinistro_data["Data do Sinistro"].strftime("%d/%m/%Y")
    elif re.match(r"\d{4}-\d{2}-\d{2}", data_str):
        try:
            data_obj = datetime.strptime(data_str, "%Y-%m-%d %H:%M:%S")
            data_formatada = data_obj.strftime("%d/%m/%Y")
        except ValueError:
            pass # Manter a string original se a conversão falhar

    # Carregar o template PPT
    template_path = 'C:\\Fabio\\Desenvolvimento\\Varejo\\ppt\\Sinistro_65329.pptx'
    output_path = f'C:\\Fabio\\Desenvolvimento\\Varejo\\ppt\\{numero_sinistro}_Gerado.pptx'
    
    # Criar uma cópia do template
    prs = Presentation(template_path)
    
    # Mapear os dados do sinistro para os campos do template
    # Usar nomes de colunas reais da planilha
    sinistro_info = {
        'Número do Sinistro': str(numero_sinistro),
        'Status': str(sinistro_data["Observação"]),
        'Data': data_formatada,
        'Causa': str(sinistro_data["Causa Final"]),
        'Transportador': str(sinistro_data["Transportador"]),
        'UF Origem': str(sinistro_data["UF - Origem"]),
        'Cidade Origem': str(sinistro_data["Cidade Origem"]),
        'UF Destino': str(sinistro_data["UF - Destino"]),
        'Cidade Destino': str(sinistro_data["Cidade - Destino"]),
        'Prejuízo Apurado': prejuizo_formatado,
        'Descrição': str(sinistro_data["Ação"])
    }
    
    # Função para substituir texto mantendo a formatação
    def replace_text_in_shape(shape, replacements):
        if not hasattr(shape, "text_frame"):
            return
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, value)
    
    # Atualizar os slides com os dados do sinistro
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                replace_text_in_shape(shape, sinistro_info)
            
            # Processar tabelas, se houver
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_shape(cell, sinistro_info)
    
    # Salvar a apresentação
    prs.save(output_path)
    print(f"Apresentação gerada com sucesso: {output_path}")
    return True

if __name__ == "__main__":
    # Se um número de sinistro for fornecido como argumento, use-o
    if len(sys.argv) > 1:
        try:
            numero_sinistro = int(sys.argv[1])
            criar_apresentacao_sinistro(numero_sinistro)
        except ValueError:
            print("Erro: O número do sinistro deve ser um valor numérico.")
    else:
        # Caso contrário, use o número de sinistro padrão (65346)
        criar_apresentacao_sinistro(65247)

