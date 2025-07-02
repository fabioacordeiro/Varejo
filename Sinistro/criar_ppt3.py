from pptx import Presentation
import re

# Caminhos dos arquivos
dados_txt_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\DADOS_SINISTRO.txt"
ppt_template_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_63589.pptx"
output_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_Sinistro_Atualizado3.pptx"


# Carregar os dados do txt
with open(dados_txt_path, "r", encoding="utf-8") as f:
    dados = f.read()

# Converter o conteúdo do txt em dicionário
dados_dict = dict(re.findall(r"(?m)^([\w\s\/\-\(\):]+):\s*(.*)", dados))

# Adicionar substituições adicionais que aparecem no slide
if "SINISTRO" in dados_dict:
    dados_dict["SINISTRO 63.589"] = f"SINISTRO {dados_dict['SINISTRO']}"

if "MOTIVO" in dados_dict:
    dados_dict["VARIAÇÃO DE TEMPERATURA"] = dados_dict["MOTIVO"]

if "Origem" in dados_dict:
    dados_dict["ITAPECERICA DA SERRA-SP"] = dados_dict["Origem"]

if "Destino" in dados_dict:
    dados_dict["4677 - PIRACICABA"] = dados_dict["Destino"]

# Carregar o template do PowerPoint
prs = Presentation(ppt_template_path)

# Substituição de texto considerando formas (shapes) com table ou text_frame
def substituir_texto_completo(prs, dados_dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for chave, valor in dados_dict.items():
                        if chave in p.text:
                            p.text = p.text.replace(chave, valor)
            elif shape.shape_type == 19:  # Tabela
                for row in shape.table.rows:
                    for cell in row.cells:
                        for chave, valor in dados_dict.items():
                            if chave in cell.text:
                                cell.text = cell.text.replace(chave, valor)

# Aplicar substituição
substituir_texto_completo(prs, dados_dict)

# Salvar apresentação atualizada
prs.save(output_path)

output_path