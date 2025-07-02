# Recarregar arquivos após reset do ambiente
from pptx import Presentation
import re

# Caminhos dos arquivos
dados_txt_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\DADOS_SINISTRO.txt"
ppt_template_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_63589.pptx"
output_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_Sinistro_Atualizado.pptx"

# Recarregar conteúdo do arquivo de texto
with open(dados_txt_path, "r", encoding="utf-8") as f:
    dados = f.read()

# Extrair os dados com regex
dados_dict = dict(re.findall(r"(?m)^([\w\s\/\-\(\):]+):\s*(.*)", dados))

# Carregar o template PPTX
prs = Presentation(ppt_template_path)

# Substituir texto em todos os slides
def substituir_textos_em_apresentacao(prs, dados_dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for chave, valor in dados_dict.items():
                            if chave in run.text:
                                run.text = run.text.replace(chave, valor)

# Adicionar mapeamentos para substituição baseada em texto original do template
if "SINISTRO" in dados_dict:
    dados_dict["SINISTRO 63.589"] = f"SINISTRO {dados_dict['SINISTRO']}"

if "MOTIVO" in dados_dict:
    dados_dict["VARIAÇÃO DE TEMPERATURA"] = dados_dict["MOTIVO"]

if "Origem" in dados_dict:
    dados_dict["ITAPECERICA DA SERRA-SP"] = dados_dict["Origem"]

if "Destino" in dados_dict:
    dados_dict["4677 - PIRACICABA"] = dados_dict["Destino"]

# Aplicar substituições
substituir_textos_em_apresentacao(prs, dados_dict)

# Salvar o novo arquivo
prs.save(output_path)

output_path