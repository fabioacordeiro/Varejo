import pandas as pd

# Carregar o arquivo CSV para visualizar os dados
csv_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\Carga.csv"
df_carga = pd.read_csv(csv_path, sep=None, engine='python')  # Detectar automaticamente o separador

# Mostrar as primeiras linhas do DataFrame
df_carga.head()

from pptx import Presentation
from pptx.util import Inches
import re

# Carregar a apresentação de template
template_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\Resumo_63589.pptx"
ppt = Presentation(template_path)

# Obter os dados do CSV
row = df_carga.iloc[0]
numero_carga = row["\ufeffNúmero da Carga"]
transportador = row["Transportador"]
filial = row["Filial"]
data_emissao_nf = row["Data da Carga"]
tipo_carga = row["Tipo de Carga"]

# Substituir textos no slide com os dados do CSV
# 3. Substituições a serem feitas
substituicoes = {
    "63.589": numero_carga,
    "TRANSPORTES FRAMENTO LTDA": transportador,
    "ITAPECERICA DA SERRA-SP": filial,
    "25/03/2025 04:43": data_emissao_nf,
    "RESFRIADO": tipo_carga,
}




# Redefinir a função para substituir textos nos slides
def substituir_texto_no_slide(slide, substituicoes):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    for chave, valor in substituicoes.items():
                        if chave in run.text:
                            run.text = run.text.replace(chave, valor)

# Aplicar as substituições no primeiro slide novamente
slide = ppt.slides[0]
substituir_texto_no_slide(slide, substituicoes)

# Salvar a nova apresentação
output_path = f"C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\Resumo_{numero_carga}.pptx"
ppt.save(output_path)

output_path

print("Concluído")