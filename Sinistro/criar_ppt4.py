#pip install pyxlsb

import pandas as pd
import pyxlsb
from pptx import Presentation
import tabula

try:
    # 1. Leitura dos dados do arquivo CSV (df_carga)
    df_carga = pd.read_csv("C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Carga.csv")

    # 2. Leitura dos dados do arquivo PDF (df_temperatura)
    df_temperatura = tabula.read_pdf("C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Temperatura.pdf", pages=1)[0]

    # 3. Leitura do arquivo XLSB usando pyxlsb e pandas
    with pyxlsb.open_workbook('C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Grafico.xlsb') as wb:
        sheets = []
        for sheet_name in wb.sheets:
            with wb.get_sheet(sheet_name) as sheet:
                data = []
                for row in sheet.rows():
                    data.append([item.v for item in row])
                sheets.append(pd.DataFrame(data))
        df_xlsb = pd.concat(sheets)

    # 4. Inserir dados do df_carga no DataFrame df_xlsb
    # Supondo que você queira inserir os dados do df_carga nas primeiras colunas do df_xlsb
    for col_num, col_name in enumerate(df_carga.columns):
        df_xlsb.insert(col_num, col_name, df_carga[col_name])

    # 5. Geração do arquivo Excel (XLSX)
    df_xlsb.to_excel("C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Grafico.xlsx", index=False)

    print("Arquivo XLSB convertido para XLSX com os dados do df_carga!")

    # 6. Geração do arquivo PowerPoint (PPTX)
    presentation = Presentation("C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_63589_template.pptx")

    # Substitui os placeholders no PowerPoint com os dados de df_carga e df_temperatura
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "carga_total" in shape.text:
                    shape.text_frame.text = str(df_carga["Carga Total"].sum())  # Exemplo: Soma da carga total
                elif "temperatura_media" in shape.text:
                    shape.text_frame.text = str(df_temperatura["Temperatura"].mean())  # Exemplo: Média da temperatura

    # Salva o arquivo PowerPoint
    presentation.save("C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_63589.pptx")

except Exception as e:
    print(f"Ocorreu um erro: {e}")