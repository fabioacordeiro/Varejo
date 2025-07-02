#
# pip install pandas openpyxl python-pptx tabula-py
# pip install JPype1

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
import tabula

# 1. Leitura dos dados do arquivo CSV
df_carga = pd.read_csv("Carga.csv")

# 2. Leitura dos dados do arquivo PDF
# Assumindo que a tabela de temperatura está na primeira página e é a primeira tabela encontrada
df_temperatura = tabula.read_pdf("Temperatura.pdf", pages=1)[0]

# 3. Manipulação dos dados (adapte conforme necessário)
# Exemplo: Combinar os dados dos dois DataFrames
df_combinado = pd.concat([df_carga, df_temperatura], axis=1)

# 4. Geração do arquivo Excel (XLSB)
# Carrega o template do Excel
workbook = load_workbook("Grafico_template.xlsb")
sheet = workbook.active

# Escreve os dados no arquivo Excel
for r in range(df_combinado.shape[0]):
    for c in range(df_combinado.shape[1]):
        sheet.cell(row=r + 2, column=c + 1).value = df_combinado.iloc[r, c]

# Salva o arquivo Excel
workbook.save("Grafico.xlsb")

# 5. Geração do arquivo PowerPoint (PPTX)
# Carrega o template do PowerPoint
presentation = Presentation("Resumo_63589_template.pptx")

# Substitui os placeholders no PowerPoint com os dados
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "carga_total" in shape.text:
                shape.text_frame.text = str(df_carga["Carga Total"].sum())  # Exemplo: Soma da carga total
            elif "temperatura_media" in shape.text:
                shape.text_frame.text = str(df_temperatura["Temperatura"].mean())  # Exemplo: Média da temperateratura
# Salva o arquivo PowerPoint
presentation.save("Resumo.pptx")