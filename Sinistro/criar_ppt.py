from pptx import Presentation
import re

# Carregar o conteúdo do arquivo de texto
dados_txt_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\DADOS_SINISTRO.txt"
with open(dados_txt_path, "r", encoding="utf-8") as f:
    dados = f.read()

# Extrair os dados com regex
dados_dict = dict(re.findall(r"(?m)^([\w\s\/\-\(\):]+):\s*(.*)", dados))

# Visualize os dados extraídos (para debug)
for k, v in dados_dict.items():
    print(f"{k}: {v}")

# Carregar o template PPTX
ppt_template_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_63589.pptx"
prs = Presentation(ppt_template_path)

# Função robusta para substituir textos em um slide
def substituir_textos_em_slide(slide, substituicoes: dict):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for placeholder, novo_texto in substituicoes.items():
                if placeholder in paragraph.text:
                    # Substitui em cada run do parágrafo
                    for run in paragraph.runs:
                        run.text = run.text.replace(placeholder, novo_texto)

# Definir dicionários de substituição para cada slide
# OBS.: Certifique-se de que os placeholders abaixo correspondem exatamente ao que está no template.
substituicoes_slide0 = {
    "SINISTRO 63.589": f"SINISTRO {dados_dict.get('SINISTRO', '')}",
    "VARIAÇÃO DE TEMPERATURA": dados_dict.get("MOTIVO", ""),
    "ITAPECERICA DA SERRA-SP": dados_dict.get("Origem", ""),
    "4677 - PIRACICABA": dados_dict.get("Destino", ""),
}

substituicoes_slide1 = {
    "ENCOSTA EM DOCA": dados_dict.get("ENCOSTA EM DOCA", ""),
    "INICIO CARREGAMENTO": dados_dict.get("INICIO CARREGAMENTO", ""),
    "FIM CARREGAMENTO": dados_dict.get("FIM CARREGAMENTO", ""),
    "EMISSÃO NF": dados_dict.get("EMISSÃO NF", ""),
    "INICIO DE VIAGEM": dados_dict.get("INICIO DE VIAGEM", ""),
    "CHEGADA EM LOJA/CD": dados_dict.get("CHEGADA EM LOJA/CD", ""),
    # Se o slide 1 também contiver o identificador do sinistro:
    "SINISTRO 63.589": f"SINISTRO {dados_dict.get('SINISTRO', '')}"
}

# Aplicar as substituições nos slides desejados
# Se você tiver mais de 2 slides, poderá criar outros dicionários conforme a necessidade.
if len(prs.slides) >= 1:
    substituir_textos_em_slide(prs.slides[0], substituicoes_slide0)
if len(prs.slides) >= 2:
    substituir_textos_em_slide(prs.slides[1], substituicoes_slide1)

# Salvar o novo arquivo
output_path = "C:\\Fabio\\Desenvolvimento\\Varejo\\Sinistro\\ppt\\Resumo_Sinistro_Atualizado.pptx"
prs.save(output_path)
print("Novo PPTX salvo em:", output_path)