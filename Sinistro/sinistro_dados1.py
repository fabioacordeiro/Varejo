import os
import pandas as pd
import tabula
from openpyxl import load_workbook
from pptx import Presentation

try:
    # 1. Leitura dos dados do arquivo CSV
    df_carga = pd.read_csv("Carga.csv")

    # 2. Leitura dos dados do arquivo PDF
    df_temperatura = tabula.read_pdf("Temperatura.pdf", pages=1)[0]

    # 3. Combinar os DataFrames usando pd.merge
    # Substitua "ColunaComum" pelo nome real da coluna em comum
    df_combinado = pd.merge(df_carga, df_temperatura, on="ColunaComum")

    # 4. Geração do arquivo Excel (XLSB)
    workbook = load_workbook("Grafico_template.xlsb")
    sheet = workbook.active

    # Escreve os dados no arquivo Excel
    for r in range(df_combinado.shape[0]):
        for c in range(df_combinado.shape[1]):
            sheet.cell(row=r + 2, column=c + 1).value = df_combinado.iloc[r, c]

    # Salva o arquivo Excel
    workbook.save("Grafico.xlsb")

    # 5. Geração do arquivo PowerPoint (PPTX)
    presentation = Presentation("Resumo_63589_template.pptx")

    # Substitui os placeholders no PowerPoint com os dados
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "carga_total" in shape.text:
                    shape.text_frame.text = str(df_carga["Carga Total"].sum())  # Exemplo: Soma da carga total
                elif "temperatura_media" in shape.text:
                    shape.text_frame.text = str(df_temperatura["Temperatura"].mean())  # Exemplo: Média da temperatura

    # Salva o arquivo PowerPoint
    presentation.save("Resumo_63589.pptx")

except Exception as e:
    print(f"Ocorreu um erro: {e}")